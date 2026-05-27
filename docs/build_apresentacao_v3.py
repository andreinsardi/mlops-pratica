"""Apresentacao v3 - usando os layouts FGV CORRETOS do template.

Layouts FGV usados:
- "Slide de abertura"  -> capa e divisores (title idx=0 + body idx=10)
- "Titulo 1 linha - conteudo"  -> bullets puro (title idx=0 + body idx=14)
- "Titulo 1 linha - conteudo + imagem_d"  -> texto+imagem (title idx=0 + body idx=19 + pic idx=12)
- "Conteudo"  -> diagrama full (sem titulo - apenas body idx=14 que vira container da imagem)
- "1_slideFinal"  -> obrigado (idx=15)

Tamanho do slide: 20 x 11.25 in (super-wide FGV).
"""
from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# -----------------------------------------------------------------------------
TEMPLATE = Path("/sessions/busy-hopeful-clarke/mnt/uploads/MLOps_E1_ E2_DevOps_MLOps.pptx")
PROJECT_DIR = Path("/sessions/busy-hopeful-clarke/mnt/FGV/mlops-pratica")
DOCS_DIR = PROJECT_DIR / "docs"
IMG_DIR = DOCS_DIR / "img"
MMD_DIR = DOCS_DIR / "mermaid"
DOT_DIR = DOCS_DIR / "dot"
OUT_PPTX = PROJECT_DIR / "Apresentacao_MLOps_Pratica.pptx"
for d in [IMG_DIR, MMD_DIR, DOT_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# Cores
AZUL = RGBColor(0x1F, 0x4E, 0x79)
VERMELHO = RGBColor(0xC0, 0x00, 0x00)
CINZA = RGBColor(0x59, 0x59, 0x59)
CINZA_ESC = RGBColor(0x3F, 0x3F, 0x3F)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

HX = {
    "azul": "#1F4E79", "verde": "#5A8F4B", "vermelho": "#C00000",
    "cinza": "#595959", "laranja": "#ED7D31", "branco": "#FFFFFF",
}


# =============================================================================
# DIAGRAMAS (Mermaid source + DOT render)
# =============================================================================
MERMAID_SRC = {
    "01_arquitetura": """flowchart LR
  API[HackerNews API] --> AF[Airflow]
  AF --> MIN[(MinIO data lake)]
  MIN --> MLF[MLflow Tracking+Registry]
  MIN --> M1[Modelo Preditivo]
  MIN --> M2[Modelo Embeddings]
  M1 --> MLF
  M2 --> MLF
  M2 --> PG[(Postgres pgvector)]
  MLF --> SRV[FastAPI /predict /search]
  PG --> SRV
""",
    "02_eng_dados": """flowchart LR
  A[extract HN] --> B[raw parquet]
  B --> C[curate dedup]
  C --> D1[features tabular]
  C --> D2[features text]
  D1 --> E1[treino preditivo]
  D2 --> E2[treino embeddings]
""",
    "03_mlops": """flowchart LR
  L[Load features] --> T[Train RF+XGBoost]
  T --> E[Evaluate]
  E --> M[MLflow Tracking]
  M --> R[Registry Staging]
  R --> P[Production]
  P --> API[FastAPI /predict]
""",
    "04_devops": """flowchart LR
  F[feature/*] --> D[develop] --> R[release/*] --> MA[main] --> T[tag v*]
  F -.-> CI1[CI ruff+pytest]
  R -.-> CI2[CI docker build]
  MA -.-> CD1[CD staging]
  T -.-> CD2[CD producao]
""",
    "05_dados": """flowchart LR
  API[HN API] --> RAW[raw parquet]
  RAW --> CUR[curated parquet]
  CUR --> FT[features tabular]
  CUR --> FX[features text]
  FT --> PG[(Postgres pgvector)]
  FX --> PG
""",
    "06_gcp": """flowchart LR
  INET[Internet] --> FW[Firewall]
  FW --> VM[Compute Engine]
  VM --> DISK[(Persistent SSD)]
  VM -.-> CSQL[(Cloud SQL futuro)]
  VM -.-> GCS[(GCS bucket futuro)]
""",
    "07_solucao_geral": """flowchart TB
  subgraph PIPE[Backend MLOps Docker]
    AF[Airflow] --> MIN[MinIO]
    AF --> MLF[MLflow]
    MIN --> MLF
  end
  subgraph CLI[Cliente Final]
    P[FastAPI /predict]
    S[FastAPI /search]
    UI[UIs Airflow/MLflow/MinIO]
  end
  MLF --> P
  MIN --> S
""",
    "08_airflow_mlflow": """flowchart LR
  subgraph AF[Airflow ORQUESTRA]
    AF1[Quando rodar]
    AF2[Em que ordem]
    AF3[Retry/alert]
  end
  subgraph ML[MLflow VERSIONA MODELOS]
    ML1[Codigo do experimento]
    ML2[Params + metrics]
    ML3[Artifacts pickle]
    ML4[Stages do Registry]
  end
  subgraph LK[MinIO VERSIONA DADOS]
    LK1[raw por hora]
    LK2[curated]
    LK3[features]
  end
""",
}


def write_mermaid_files():
    for name, src in MERMAID_SRC.items():
        (MMD_DIR / f"{name}.mmd").write_text(src, encoding="utf-8")


# Aspect ratio dos diagramas - todos LR (horizontais)
def dot_lr_header():
    return ("""digraph G {
    rankdir=LR;
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=14,
          fontcolor="white", margin="0.25,0.15", penwidth=0];
    edge [color="#595959", penwidth=2, arrowsize=1.0];
    splines=spline;
    nodesep=0.6; ranksep=1.0;
""")


def dot_node(nid, label, fill="#1F4E79", shape="box", fontcolor="white"):
    label = label.replace("\n", "\\n")
    return f'    {nid} [label="{label}", fillcolor="{fill}", shape={shape}, fontcolor="{fontcolor}"];\n'


def build_dot_diagrams():
    diags = {}

    # 01 arquitetura
    d = dot_lr_header()
    d += dot_node("API", "HackerNews\\nAPI", HX["verde"])
    d += dot_node("AF", "Airflow\\nweb + scheduler", HX["azul"])
    d += dot_node("MIN", "MinIO\\ndata lake S3", HX["azul"], "cylinder")
    d += dot_node("MLF", "MLflow\\nTracking + Registry", HX["azul"])
    d += dot_node("M1", "ML Preditivo\\nsklearn + XGBoost", HX["cinza"])
    d += dot_node("M2", "Embeddings\\nMiniLM-L6 (CPU)", HX["cinza"])
    d += dot_node("PG", "Postgres + pgvector", HX["laranja"], "cylinder")
    d += dot_node("SRV", "FastAPI\\n/predict /search", HX["azul"])
    for a, b in [("API","AF"),("AF","MIN"),("MIN","MLF"),
                 ("MIN","M1"),("MIN","M2"),("M1","MLF"),
                 ("M2","MLF"),("M2","PG"),("MLF","SRV"),("PG","SRV")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["01_arquitetura"] = d

    # 02 eng dados
    d = dot_lr_header()
    d += dot_node("E", "1. extract\\nHN topstories", HX["verde"])
    d += dot_node("R", "2. raw\\nparquet", HX["azul"])
    d += dot_node("C", "3. curate\\ndedup por id", HX["azul"])
    d += dot_node("FT", "4a. features\\ntabular", HX["laranja"])
    d += dot_node("FX", "4b. features\\ntext", HX["laranja"])
    d += dot_node("TP", "treino\\npreditivo", HX["cinza"])
    d += dot_node("TE", "treino\\nembeddings", HX["cinza"])
    for a, b in [("E","R"),("R","C"),("C","FT"),("C","FX"),
                 ("FT","TP"),("FX","TE")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["02_eng_dados"] = d

    # 03 mlops
    d = dot_lr_header()
    d += dot_node("L", "Load features\\nMinIO", HX["azul"])
    d += dot_node("T", "Train\\nRF + XGBoost", HX["azul"])
    d += dot_node("EV", "Evaluate\\nAUC F1 PR-AUC", HX["cinza"])
    d += dot_node("M", "MLflow Tracking\\nparams metrics\\nartifacts signature", HX["laranja"])
    d += dot_node("R", "Registry\\nNone -> Staging", HX["azul"])
    d += dot_node("P", "Promote\\nProduction", HX["verde"])
    d += dot_node("API", "FastAPI\\n/predict /search", HX["azul"])
    for a, b in [("L","T"),("T","EV"),("EV","M"),("M","R"),("R","P"),("P","API")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["03_mlops"] = d

    # 04 devops
    d = dot_lr_header()
    d += dot_node("F", "feature/*", HX["verde"])
    d += dot_node("D", "develop", HX["laranja"])
    d += dot_node("R", "release/*", HX["cinza"])
    d += dot_node("MA", "main", HX["azul"])
    d += dot_node("T", "tag v*", HX["azul"])
    d += dot_node("CI1", "CI feature\\nruff + pytest", HX["verde"])
    d += dot_node("CI2", "CI release\\ndocker build", HX["cinza"])
    d += dot_node("CD1", "CD staging", HX["azul"])
    d += dot_node("CD2", "CD producao", HX["laranja"])
    for a, b in [("F","D"),("D","R"),("R","MA"),("MA","T")]:
        d += f"    {a} -> {b};\n"
    for a, b in [("F","CI1"),("R","CI2"),("MA","CD1"),("T","CD2")]:
        d += f'    {a} -> {b} [style="dashed"];\n'
    d += "}"
    diags["04_devops"] = d

    # 05 dados
    d = dot_lr_header()
    d += dot_node("API", "HN API\\nfirebaseio", HX["verde"])
    d += dot_node("RAW", "raw parquet\\nMinIO\\ndt=YYYY-MM-DD/hr=HH", HX["azul"])
    d += dot_node("CUR", "curated\\ndedup por id", HX["azul"])
    d += dot_node("FT", "features tabular", HX["laranja"])
    d += dot_node("FX", "features text", HX["laranja"])
    d += dot_node("PG", "Postgres + pgvector\\nairflow / mlflow / app", HX["laranja"], "cylinder")
    for a, b in [("API","RAW"),("RAW","CUR"),("CUR","FT"),
                 ("CUR","FX"),("FT","PG"),("FX","PG")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["05_dados"] = d

    # 06 GCP
    d = dot_lr_header()
    d += dot_node("INET", "Internet", HX["cinza"])
    d += dot_node("FW", "Firewall\\ntag mlops\\n8080/5000/9001/8000", HX["laranja"])
    d += dot_node("VM", "Compute Engine\\ne2-standard-4\\nUbuntu 22.04", HX["azul"])
    d += dot_node("DISK", "Persistent SSD\\n200GB", HX["cinza"], "cylinder")
    d += dot_node("CSQL", "Cloud SQL\\n(evolucao)", HX["verde"], "cylinder")
    d += dot_node("GCS", "GCS bucket\\n(evolucao)", HX["verde"], "cylinder")
    for a, b in [("INET","FW"),("FW","VM"),("VM","DISK")]:
        d += f"    {a} -> {b};\n"
    for a, b in [("VM","CSQL"),("VM","GCS")]:
        d += f'    {a} -> {b} [style="dashed"];\n'
    d += "}"
    diags["06_gcp"] = d

    # 07 solucao - layout horizontal com clusters
    d = """digraph G {
    rankdir=LR;
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=14,
          fontcolor="white", margin="0.25,0.15", penwidth=0];
    edge [color="#595959", penwidth=2, arrowsize=1.0];
    nodesep=0.5; ranksep=1.2;

    subgraph cluster_pipe {
        label="Backend MLOps (Docker - interno)";
        labeljust=l; labelloc=t;
        style="rounded,filled"; fillcolor="#EAF3F8"; color="#1F4E79";
        fontcolor="#1F4E79"; fontsize=16; fontname="Helvetica";
        B1 [label="Airflow\\n3 DAGs", fillcolor="#1F4E79"];
        B2 [label="MLflow\\nTracking + Registry", fillcolor="#1F4E79"];
        B3 [label="MinIO\\ndata lake", fillcolor="#1F4E79"];
        B4 [label="Postgres\\npgvector", fillcolor="#ED7D31"];
    }
    subgraph cluster_client {
        label="Cliente Final consome";
        labeljust=l; labelloc=t;
        style="rounded,filled"; fillcolor="#E2F0D9"; color="#5A8F4B";
        fontcolor="#5A8F4B"; fontsize=16; fontname="Helvetica";
        A1 [label="FastAPI\\n/predict", fillcolor="#5A8F4B"];
        A2 [label="FastAPI\\n/search semantico", fillcolor="#5A8F4B"];
        A3 [label="UIs operacionais\\nAirflow MLflow MinIO", fillcolor="#5A8F4B"];
    }
    B1 -> B3;
    B1 -> B2;
    B2 -> B3;
    B2 -> A1;
    B3 -> A1;
    B4 -> A2;
    B3 -> A3;
}
"""
    diags["07_solucao_geral"] = d

    # 08 responsabilidades - 3 clusters verticais
    d = """digraph G {
    rankdir=TB;
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=13,
          fontcolor="white", margin="0.20,0.12", penwidth=0];
    edge [style=invis];
    nodesep=0.3; ranksep=0.4;

    subgraph cluster_af {
        label="Airflow ORQUESTRA";
        labeljust=l; labelloc=t; rank=same;
        style="rounded,filled"; fillcolor="#EAF3F8"; color="#1F4E79";
        fontcolor="#1F4E79"; fontsize=18; fontname="Helvetica";
        AF1 [label="Quando rodar\\n(cron schedule)", fillcolor="#1F4E79"];
        AF2 [label="Em que ordem\\n(dependencias)", fillcolor="#1F4E79"];
        AF3 [label="Retry e alerta\\n(falhas)", fillcolor="#1F4E79"];
        AF4 [label="Idempotencia\\n(reexecutar safe)", fillcolor="#1F4E79"];
        AF1 -> AF2 -> AF3 -> AF4;
    }
    subgraph cluster_ml {
        label="MLflow VERSIONA MODELOS";
        labeljust=l; labelloc=t; rank=same;
        style="rounded,filled"; fillcolor="#FFF2CC"; color="#C28A00";
        fontcolor="#C28A00"; fontsize=18; fontname="Helvetica";
        ML1 [label="Codigo do experimento\\n(run id)", fillcolor="#ED7D31"];
        ML2 [label="Params + metrics\\n(track)", fillcolor="#ED7D31"];
        ML3 [label="Artifacts\\n(pipeline pickle)", fillcolor="#ED7D31"];
        ML4 [label="Stages do Registry\\n(None/Staging/Prod)", fillcolor="#ED7D31"];
        ML1 -> ML2 -> ML3 -> ML4;
    }
    subgraph cluster_lk {
        label="MinIO VERSIONA DADOS";
        labeljust=l; labelloc=t; rank=same;
        style="rounded,filled"; fillcolor="#E2F0D9"; color="#5A8F4B";
        fontcolor="#5A8F4B"; fontsize=18; fontname="Helvetica";
        LK1 [label="raw por hora\\n(imutavel)", fillcolor="#5A8F4B"];
        LK2 [label="curated\\n(dedup snapshot)", fillcolor="#5A8F4B"];
        LK3 [label="features\\n(materializado)", fillcolor="#5A8F4B"];
        LK1 -> LK2 -> LK3;
    }
}
"""
    diags["08_airflow_mlflow"] = d

    return diags


def render_dot_to_png(dots):
    images = {}
    for name, content in dots.items():
        dot_path = DOT_DIR / f"{name}.dot"
        png_path = IMG_DIR / f"{name}.png"
        dot_path.write_text(content, encoding="utf-8")
        result = subprocess.run(
            ["dot", "-Tpng", "-Gdpi=160", str(dot_path), "-o", str(png_path)],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(f"dot falhou para {name}: {result.stderr}")
        images[name] = str(png_path)
    return images


# =============================================================================
# PPTX usando layouts FGV reais
# =============================================================================
def clone_template():
    prs = Presentation(str(TEMPLATE))
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.attrib[qn("r:id")]
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)
    return prs


def get_layout(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise ValueError(f"Layout '{name}' nao encontrado")


def _run_set(run, text, size=None, bold=None, color=None, italic=False, name="Calibri"):
    run.text = text
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    # Sempre setar bold explicitamente (default False) para evitar
    # herdar bold do master template
    run.font.bold = bool(bold) if bold is not None else False
    if italic:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def set_title(slide, text, size=36, color=AZUL):
    """Preenche o placeholder de titulo (idx=0)."""
    ph = slide.placeholders[0]
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, text, size=size, bold=True, color=color)


def fill_body_bullets(ph, lines, base_size=16):
    """Preenche um placeholder de body com bullets formatados.
    Linhas '# ...' viram subtitulos (azul, bold, sem bullet)
    Linhas vazias = espacador
    Demais = bullets normais
    """
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        # remove bullet do tema (definindo nivel a mao)
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        # buNone (sem bullet) por padrao - depois aplicamos bullet manual
        buNone = pPr.find(_qn("a:buNone"))
        if buNone is None:
            from lxml import etree
            for tag in ("a:buChar", "a:buAutoNum"):
                old = pPr.find(_qn(tag))
                if old is not None:
                    pPr.remove(old)
            pPr.append(etree.SubElement(pPr, _qn("a:buNone")))

        r = p.add_run()
        if line.startswith("# "):
            _run_set(r, line[2:], size=base_size + 4, bold=True, color=AZUL)
            p.space_before = Pt(6)
            p.space_after = Pt(4)
        elif line == "":
            _run_set(r, " ", size=8)
            p.space_after = Pt(2)
        else:
            _run_set(r, "• " + line, size=base_size, color=CINZA_ESC)
            p.space_after = Pt(4)


# -----------------------------------------------------------------------------
# Tipos de slide
# -----------------------------------------------------------------------------
def slide_capa(prs):
    """Capa - usa 'Slide de abertura' (title idx=0, body idx=10)"""
    s = prs.slides.add_slide(get_layout(prs, "Slide de abertura"))
    # Title
    ph0 = s.placeholders[0]
    tf = ph0.text_frame; tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, "Pratica MLOps end-to-end", size=44, bold=True, color=AZUL)
    # Body subtitle
    ph10 = s.placeholders[10]
    tf = ph10.text_frame; tf.clear()
    lines = [
        ("HackerNews + Airflow + MLflow + pgvector + FastAPI", 22, True, AZUL),
        ("100% Docker  |  Deploy local e em VM GCP", 18, False, CINZA),
        ("", 8, False, CINZA),
        ("Prof. Andre Insardi", 16, True, CINZA_ESC),
        ("ext.andre.insardi@prof.fgv.edu.br  |  Maio/2026", 14, False, CINZA),
    ]
    for i, (txt, sz, bd, cl) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        _run_set(r, txt, size=sz, bold=bd, color=cl)
    return s


def slide_bullets(prs, titulo, lines):
    """Bullets puros - usa 'Titulo 1 linha - conteudo' (title idx=0, body idx=14)"""
    s = prs.slides.add_slide(get_layout(prs, "Título 1 linha - conteúdo"))
    set_title(s, titulo, size=32)
    fill_body_bullets(s.placeholders[14], lines, base_size=16)
    return s


def slide_texto_imagem(prs, titulo, lines, img_path):
    """Texto a esquerda + imagem a direita - usa 'Titulo 1 linha - conteudo + imagem_d'
       (title idx=0, body idx=19, pic idx=12)"""
    s = prs.slides.add_slide(get_layout(prs, "Título 1 linha - conteúdo + imagem_d"))
    set_title(s, titulo, size=30)
    fill_body_bullets(s.placeholders[19], lines, base_size=14)
    # Inserir imagem no placeholder de imagem (idx=12)
    pic_ph = s.placeholders[12]
    pic_ph.insert_picture(img_path)
    return s


def slide_imagem_full(prs, titulo, img_path, legenda=None):
    """Titulo no topo + imagem grande ocupando o corpo.
    Usa 'Titulo 1 linha - conteudo' mas substitui o body por imagem.
    Body box: L=1.57 T=2.55 W=16.85 H=7.95
    """
    s = prs.slides.add_slide(get_layout(prs, "Título 1 linha - conteúdo"))
    set_title(s, titulo, size=30)
    # Remove o body placeholder
    body = s.placeholders[14]
    sp_elem = body._element
    sp_elem.getparent().remove(sp_elem)
    # Adiciona imagem dentro do espaco do body, centralizada
    body_L = Inches(1.57)
    body_T = Inches(2.55)
    body_W = Inches(16.85)
    body_H = Inches(7.5)  # deixa espaco para legenda embaixo
    # Adiciona imagem com largura quase total mantendo proporcao
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    # tenta fazer caber em body_W x body_H
    max_w = body_W
    target_h = Emu(int(max_w / aspect))
    if target_h > body_H:
        target_h = body_H
        max_w = Emu(int(body_H * aspect))
    # centraliza horizontalmente
    pic_L = body_L + Emu((body_W - max_w) // 2)
    pic_T = body_T + Emu((body_H - target_h) // 2)
    s.shapes.add_picture(img_path, pic_L, pic_T, width=max_w, height=target_h)

    if legenda:
        tb = s.shapes.add_textbox(body_L, Inches(10.3), body_W, Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        _run_set(r, legenda, size=12, italic=True, color=CINZA)
    return s


def slide_divisor(prs, numero, titulo, descricao):
    """Slide divisor de bloco - design com numero GRANDE a esquerda.
    Usa layout 'Slide de abertura' (mantem branding FGV) mas custom textboxes."""
    s = prs.slides.add_slide(get_layout(prs, "Slide de abertura"))
    # Apaga placeholders default
    from pptx.oxml.ns import qn as _qn
    for ph in list(s.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Numero gigante a esquerda (mas dentro da area util, depois da faixa azul)
    from pptx.enum.shapes import MSO_SHAPE
    # "BLOCO" label pequeno em cima
    tb = s.shapes.add_textbox(Inches(2.3), Inches(3.5), Inches(5.0), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, "BLOCO", size=22, bold=True, color=CINZA)
    # Numero gigante
    tb = s.shapes.add_textbox(Inches(2.3), Inches(3.9), Inches(5.0), Inches(3.8))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, numero, size=240, bold=True, color=AZUL)
    # Titulo grande a direita
    tb = s.shapes.add_textbox(Inches(8.5), Inches(4.5), Inches(10.5), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, titulo, size=48, bold=True, color=AZUL)
    # Linha vermelha separadora
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(6.7),
                              Inches(3.0), Inches(0.06))
    rect.fill.solid(); rect.fill.fore_color.rgb = VERMELHO
    rect.line.fill.background()
    # Descricao
    tb = s.shapes.add_textbox(Inches(8.5), Inches(7.0), Inches(10.5), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    _run_set(r, descricao, size=20, color=CINZA_ESC)
    return s


def slide_final(prs):
    """Slide final - usa '1_slideFinal' (idx=15 com citacao).
    A area do placeholder idx=15 fica na parte branca do slide,
    entao usamos cor AZUL para o texto ficar legivel."""
    s = prs.slides.add_slide(get_layout(prs, "1_slideFinal"))
    try:
        ph = s.placeholders[15]
        tf = ph.text_frame; tf.clear()
        tf.word_wrap = True
        lines = [
            ("Obrigado!", 32, True, AZUL),
            ("", 6, False, AZUL),
            ("Prof. Andre Insardi", 14, True, CINZA_ESC),
            ("ext.andre.insardi@prof.fgv.edu.br", 12, False, CINZA),
            ("FGV — MBA IA & Analytics — Maio/2026", 11, False, CINZA),
        ]
        for i, (t, sz, bd, cl) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            _run_set(r, t, size=sz, bold=bd, color=cl)
    except Exception as e:
        print(f"AVISO no slide final: {e}")
    return s


# =============================================================================
def main():
    print("==> Salvando .mmd")
    write_mermaid_files()
    print("==> Renderizando PNGs via Graphviz")
    dots = build_dot_diagrams()
    imgs = render_dot_to_png(dots)
    print("==> Construindo PPTX com layouts FGV")
    prs = clone_template()

    # 1 - Capa
    slide_capa(prs)

    # 2 - Agenda
    slide_bullets(prs, "Agenda da Pratica Integradora", [
        "# Roteiro (9 blocos)",
        "0. Visao da solucao - o que sera entregue ao cliente",
        "1. Problema e caso de uso (HN 'vai bombar')",
        "2. Base de dados (API + camadas + Postgres/pgvector)",
        "3. Arquitetura geral (7 containers Docker)",
        "4. Engenharia de Dados (Airflow)",
        "5. Fluxo MLOps (MLflow Tracking + Registry + Serving)",
        "6. Fluxo DevOps (Gitflow + CI/CD)",
        "7. Deploy em VM GCP",
        "8. Conclusoes e proximos passos",
    ])

    # ===== BLOCO 0 - Visao da Solucao =====
    slide_divisor(prs, "0", "Visao da Solucao",
        "O que sera entregue ao cliente, o que o Airflow controla, o que o MLflow versiona e como os dados sao versionados.")
    slide_bullets(prs, "O que vamos construir (resumo executivo)", [
        "# Visao em uma frase",
        "Plataforma MLOps rodando em Docker que ingere dados publicos do HackerNews, treina 2 pipes de ML (preditivo + embeddings) e expoe servicos via API REST.",
        "",
        "# Componentes (7 containers)",
        "Postgres + pgvector (metadata + vector store)",
        "MinIO (data lake S3-compatible)",
        "MLflow (Tracking + Registry)",
        "Airflow webserver + scheduler (orquestra 3 DAGs)",
        "FastAPI (serving /predict + /search)",
        "",
        "# Tres pipelines orquestrados",
        "Ingestao horaria de stories do HN",
        "Treino diario do classificador 'vai bombar' (RF + XGBoost)",
        "Treino diario de embeddings de titulos (MiniLM-L6) com indexacao em pgvector",
    ])
    slide_imagem_full(prs, "Visao end-to-end - Backend e Cliente Final",
                     imgs["07_solucao_geral"],
                     "Backend MLOps (azul) e oculto ao cliente; o que ele consome sao APIs e UIs (verde).")
    slide_bullets(prs, "O que o Cliente Final recebe", [
        "# Servicos de uso direto (API REST)",
        "POST /predict - dado um post novo, devolve probabilidade de 'vai bombar' (score >= 100).",
        "POST /search - busca semantica top-K por similaridade cosseno em pgvector.",
        "GET /health - liveness probe.",
        "",
        "# Interfaces operacionais",
        "Airflow UI (8080) - status das DAGs, logs, retry manual",
        "MLflow UI (5000) - experimentos, runs e Model Registry",
        "MinIO Console (9001) - exploracao do data lake",
        "Swagger UI (8000/docs) - documentacao interativa da API",
        "",
        "# SLAs implicitos (didaticos)",
        "Latencia /predict < 100ms (modelo carregado em memoria)",
        "Latencia /search < 200ms (top-10 em ate ~1M vetores no pgvector)",
        "Freshness: dados a cada hora, modelo retreinado diariamente",
    ])
    slide_imagem_full(prs, "Quem controla o que - Airflow vs MLflow vs MinIO",
                     imgs["08_airflow_mlflow"],
                     "Airflow orquestra o quando/ordem. MLflow versiona modelos. MinIO versiona dados.")
    slide_bullets(prs, "O que o Airflow controla (orquestracao)", [
        "# Responsabilidades exclusivas do Airflow",
        "QUANDO cada pipeline roda (cron schedule).",
        "EM QUE ORDEM as tasks executam (dependencias declarativas no DAG).",
        "EM CASO DE ERRO - retry com backoff exponencial e alertas.",
        "IDEMPOTENCIA - cada task pode ser re-executada sem corromper estado.",
        "OBSERVABILIDADE - logs, tempo de execucao e historico de runs.",
        "",
        "# Schedule das 3 DAGs no projeto",
        "pipeline_ingestao @hourly - extract + raw + curated + features",
        "pipeline_treino_preditivo 02:00 daily - train RF/XGBoost + register",
        "pipeline_treino_embeddings 03:00 daily - encode MiniLM + upsert pgvector",
        "",
        "# O que o Airflow NAO faz",
        "Nao versiona codigo (Git).",
        "Nao versiona modelos (MLflow Registry).",
        "Nao armazena features (MinIO).",
        "Nao serve predicoes (FastAPI).",
    ])
    slide_bullets(prs, "O que o MLflow versiona (modelos)", [
        "# Tracking (durante o treino)",
        "params - hiperparametros do estimator",
        "metrics - accuracy, F1, ROC-AUC, PR-AUC em test split estratificado",
        "artifacts - pipeline sklearn completo (preprocessor + estimator)",
        "signature - schema de input/output inferido automaticamente",
        "tags - modelo, dataset_version, autor, ambiente",
        "",
        "# Model Registry (versoes e stages)",
        "Cada run que registra um modelo cria nova VERSAO em hn_classifier",
        "Stages: None -> Staging -> Production -> Archived",
        "FastAPI carrega: models:/hn_classifier/Production",
        "",
        "# Dois experimentos separados",
        "hn_classifier - runs do pipe preditivo (RF + XGBoost)",
        "hn_embeddings - runs do encoder MiniLM (versao, dim, amostras)",
    ])
    slide_bullets(prs, "O que o MinIO versiona (dados)", [
        "# Tres camadas do lake",
        "raw - cada DAG horaria grava particao imutavel (dt=YYYY-MM-DD/hr=HH)",
        "curated - snapshot estavel deduplicado por id (maior score visto)",
        "features - dataset materializado para treino (tabular e text)",
        "",
        "# Como a versao acontece",
        "raw: nunca sobrescreve - cada hora gera arquivo NOVO (auditavel)",
        "curated: append-merge com dedup - dado canonico",
        "features: cada DAG de treino le a versao mais recente",
        "",
        "# Caminhos completos",
        "s3://raw/hn/dt=2026-05-26/hr=14/items.parquet",
        "s3://curated/hn/stories.parquet",
        "s3://features/tabular/train.parquet",
        "s3://features/text/titles.parquet",
        "",
        "# Evolucao para DVC ou Delta Lake (futuro)",
        "Para versionamento explicito com hash + lineage.",
    ])
    slide_bullets(prs, "Como Airflow + MLflow + MinIO conversam em runtime", [
        "# Fluxo de uma execucao de treino preditivo",
        "1. Scheduler Airflow dispara DAG no horario (02:00)",
        "2. Task load_features le features/tabular/train.parquet do MinIO via s3fs",
        "3. Task train chama mlflow.set_experiment + start_run",
        "4. MLflow loga params/metrics no Postgres (backend store)",
        "5. mlflow.sklearn.log_model envia pipeline para MinIO via S3 API (artifact store)",
        "6. mlflow.register_model cria versao N de hn_classifier",
        "7. Task promote transiciona via REST API: stage=Staging",
        "",
        "# Fluxo de uma predicao em runtime",
        "1. Cliente POST /predict no FastAPI",
        "2. FastAPI consulta MLflow: models:/hn_classifier/Production",
        "3. Lookup retorna URI s3://mlflow-artifacts/.../model",
        "4. FastAPI baixa pickle do MinIO, carrega em memoria (cache startup)",
        "5. Aplica pipeline e devolve probabilidade (< 100ms)",
    ])

    # ===== BLOCO 1 - Problema =====
    slide_divisor(prs, "1", "Problema",
        "Contexto pedagogico, gap dos encontros E1-E4 e o caso de uso escolhido.")
    slide_bullets(prs, "Contexto pedagogico - por que essa pratica existe", [
        "# Diagnostico nos encontros E1-E4",
        "E1 - fundamentos (DevOps -> MLOps, divida tecnica, ciclo de vida)",
        "E2 - experimentacao (autolog, Registry, HPO, Model Card)",
        "E3 - producao (deploy patterns, FastAPI/Docker, CI/CD, drift)",
        "E4 - governanca (LGPD, AI Act, LLMOps, RAG)",
        "",
        "# Gap identificado",
        "Encontros isolados. Falta INTEGRAR.",
        "Alunos precisam VER o ciclo completo executando.",
        "Querem entender RESPONSABILIDADES (quem orquestra, quem versiona, quem serve).",
        "",
        "# O que esta pratica entrega ao aluno",
        "Repositorio rodavel no laptop com 'make up'",
        "Mesma stack subindo em VM GCP",
        "Codigo Python modular e documentado",
        "Apresentacao narrando cada bloco com diagramas",
    ])
    slide_bullets(prs, "Caso de uso - HackerNews 'vai bombar'", [
        "# Definicao do target",
        "vai_bombar = (score >= 100), threshold didatico (~10-20% positivos)",
        "Binario escolhido por simplicidade e metricas intuitivas (AUC, F1)",
        "",
        "# Por que HackerNews",
        "API publica sem autenticacao - zero atrito de credenciais",
        "Dado heterogeneo - tabular (score, hora) + texto (titulo, URL)",
        "Volume controlavel - top 100 stories por extracao",
        "Caso analogo ao real - priorizar feed editorial, alerta de viral",
        "",
        "# Features que vamos usar",
        "Tabular: title_len, n_words_title, hour, weekday, has_url, has_question, domain, by_author",
        "Vetorial: embedding 384-d via sentence-transformers/all-MiniLM-L6-v2",
        "",
        "# Aplicacoes empresariais analogas",
        "Triagem de tickets, priorizacao de leads, churn risk, deduplicacao de topicos.",
    ])

    # ===== BLOCO 2 - Base de Dados =====
    slide_divisor(prs, "2", "Base de dados",
        "Fonte, camadas do lake e modelagem.")
    slide_bullets(prs, "Fonte e camadas de dados", [
        "# Fonte",
        "HackerNews Firebase API (publica, sem auth)",
        "Endpoints: /v0/topstories.json (IDs) e /v0/item/{id}.json (detalhe)",
        "",
        "# Lake (MinIO, formato parquet)",
        "raw/hn/dt=YYYY-MM-DD/hr=HH/items.parquet (Hive-style)",
        "curated/hn/stories.parquet (dedup por id, score maximo)",
        "features/tabular/train.parquet",
        "features/text/titles.parquet",
        "",
        "# Postgres (3 databases)",
        "airflow - metadata do scheduler",
        "mlflow - backend store do Tracking Server",
        "app - tabela embeddings(id, title, vector(384), score, ts) + indice HNSW",
    ])
    slide_imagem_full(prs, "Diagrama das camadas de dados", imgs["05_dados"],
                     "raw -> curated -> features no MinIO + Postgres com pgvector.")

    # ===== BLOCO 3 - Arquitetura =====
    slide_divisor(prs, "3", "Arquitetura geral",
        "Sete containers Docker em uma rede unica.")
    slide_bullets(prs, "Stack confirmada", [
        "# Camadas e tecnologias",
        "Orquestracao - Apache Airflow 2.9 (LocalExecutor)",
        "Tracking + Registry - MLflow 2.16",
        "Data Lake - MinIO (S3-compatible)",
        "Metadata + Vector - Postgres 16 + pgvector 0.7",
        "Embeddings - sentence-transformers all-MiniLM-L6-v2 (offline, CPU, dim=384)",
        "ML Preditivo - scikit-learn 1.5 + XGBoost 2.1",
        "Serving - FastAPI 0.115 + Uvicorn",
        "CI/CD - GitHub Actions (ruff + pytest + docker build)",
    ])
    slide_imagem_full(prs, "Arquitetura geral - 7 containers em rede unica", imgs["01_arquitetura"],
                     "Postgres + MinIO + MLflow + Airflow (web+sched) + FastAPI.")

    # ===== BLOCO 4 - Eng Dados =====
    slide_divisor(prs, "4", "Engenharia de Dados",
        "DAGs Airflow, particionamento, idempotencia.")
    slide_imagem_full(prs, "DAG pipeline_ingestao (@hourly)", imgs["02_eng_dados"],
                     "extract -> raw -> curated -> features (tabular + text).")
    slide_bullets(prs, "Padroes de design das DAGs", [
        "# TaskFlow API",
        "Decorators @dag e @task (Python puro)",
        "XCom automatico via return + assinatura entre tasks",
        "",
        "# Idempotencia e particionamento",
        "Cada execucao sobrescreve a particao da hora corrente",
        "Particionamento Hive-style (dt=, hr=)",
        "",
        "# Dedup na curated",
        "Concat (curated + raw novo) -> sort score desc -> drop_duplicates(id)",
        "",
        "# Resilience",
        "Cliente HTTP com tenacity (3 retries, backoff exponencial)",
        "ThreadPoolExecutor 16 workers para hidratar 100 items em ~5s",
        "Retry da task (2x, 2min) cobre falhas transientes da API",
    ])

    # ===== BLOCO 5 - MLOps =====
    slide_divisor(prs, "5", "Fluxo MLOps (MLflow)",
        "Tracking, Registry, promocao e serving.")
    slide_imagem_full(prs, "Fluxo MLOps completo", imgs["03_mlops"],
                     "Dois experimentos -> Registry -> FastAPI /predict.")
    slide_bullets(prs, "Triplo versionamento", [
        "# 1. Codigo - Git",
        "Branches main (estavel), develop (integracao), feature/* e release/*",
        "Tags v* para releases (v0.1, v0.2, ...)",
        "",
        "# 2. Dados - Parquet particionado em MinIO",
        "raw/ guarda o estado bruto da hora (auditavel)",
        "curated/ e a versao estavel (snapshot)",
        "features/ versiona o dataset de treino",
        "",
        "# 3. Modelos - MLflow Tracking + Registry",
        "Run captura params, metrics, artifacts e signature",
        "Versoes numericas + stages (None | Staging | Production | Archived)",
    ])
    slide_bullets(prs, "DAG treino preditivo (02:00 daily)", [
        "# Etapas",
        "load_features -> train_and_register_models -> promote",
        "",
        "# Modelos comparados",
        "RandomForest (n_estimators=200, max_depth=12)",
        "XGBoost (n_estimators=300, max_depth=6, learning_rate=0.1)",
        "Vencedor escolhido por ROC-AUC",
        "",
        "# Politica de promocao",
        "DAG promove ultima versao para Staging automaticamente",
        "Producao real: gate humano + comparativo de metricas no Registry",
    ])
    slide_bullets(prs, "DAG treino embeddings (03:00 daily)", [
        "# Etapas",
        "load_text -> encode -> index",
        "",
        "# Encoder",
        "sentence-transformers all-MiniLM-L6-v2 (offline, CPU, dim=384)",
        "Pre-baixado no Dockerfile (evita download em runtime)",
        "Normalizacao L2 dos vetores (cosine distance correto)",
        "",
        "# Indexacao",
        "Upsert em embeddings(id, title, vector, score, by_author, ts)",
        "Indice HNSW vector_cosine_ops",
    ])

    # ===== BLOCO 6 - DevOps =====
    slide_divisor(prs, "6", "Fluxo DevOps (Gitflow)",
        "Branches, CI no GitHub Actions e CD por estagio.")
    slide_imagem_full(prs, "Gitflow + CI/CD", imgs["04_devops"],
                     "ruff + pytest em PRs, docker build em release, deploy gated.")
    slide_bullets(prs, "Workflow ci.yml (GitHub Actions)", [
        "# Triggers",
        "push em main/develop e pull_request em main/develop",
        "",
        "# Job 1 - lint-test",
        "Setup Python 3.11 + pip cache",
        "ruff check src tests dags",
        "pytest --cov=src/mlops_pratica --cov-report=term-missing",
        "",
        "# Job 2 - docker-build (needs lint-test)",
        "docker buildx setup",
        "Build sanity-check das 3 imagens: airflow, mlflow, fastapi",
        "",
        "# Politica de release",
        "Tag vX.Y na main dispara workflow de release",
        "Gate humano antes de qualquer deploy compartilhado",
    ])

    # ===== BLOCO 7 - Deploy GCP =====
    slide_divisor(prs, "7", "Deploy em VM GCP",
        "Compute Engine + Persistent Disk + Firewall via script gcloud.")
    slide_imagem_full(prs, "Arquitetura GCP (VM unica)", imgs["06_gcp"],
                     "Firewall por tag mlops libera as portas das UIs.")
    slide_bullets(prs, "Como provisionar", [
        "# Pre-requisitos",
        "Projeto GCP com billing",
        "gcloud autenticado e cota para 4 vCPU",
        "",
        "# Comando unico",
        "export GCP_PROJECT=meu-projeto",
        "./infra/gcp/provision.sh",
        "",
        "# O script faz",
        "Cria disco SSD 200GB (pd-ssd)",
        "Cria VM ubuntu-22 com startup-script",
        "Abre firewall por tag (ajuste SOURCE_IP_RANGE em prod!)",
        "",
        "# Startup-script automatiza",
        "Instala Docker Engine + Compose plugin",
        "Formata e monta disco em /opt/mlops",
        "Clona repositorio e roda docker compose up -d",
        "",
        "# Custo estimado",
        "~US$ 150/mes (VM 24/7 + 300GB SSD). Preemptible reduz ~70%.",
    ])

    # ===== BLOCO 8 - Conclusoes =====
    slide_divisor(prs, "8", "Conclusoes",
        "O que entregamos e proximos passos para evolucao.")
    slide_bullets(prs, "O que entregamos", [
        "# Componentes implementados",
        "7 containers Docker orquestrados via docker-compose",
        "3 DAGs Airflow (ingestao, treino preditivo, treino embeddings)",
        "MLflow Tracking + Registry com promocao automatizada",
        "Vector search em pgvector com indice HNSW cosine",
        "FastAPI servindo /predict (Registry) e /search (vetorial)",
        "Triplo versionamento: Git + lake parquet + MLflow Registry",
        "CI completo (ruff + pytest + docker build) via GitHub Actions",
        "Manual de deploy em VM GCP com script de provisionamento",
    ])
    slide_bullets(prs, "Proximos passos (evolucao)", [
        "# Curto prazo (laboratorio)",
        "Monitoracao de drift com Evidently AI (KS, PSI, Wasserstein)",
        "Datasheet for Datasets + Model Card por modelo registrado",
        "Testes E2E via docker compose + pytest-compose",
        "Avaliacao de fairness (Disparate Impact, Equal Opportunity)",
        "",
        "# Medio prazo (cloud)",
        "Migrar Postgres para Cloud SQL",
        "Migrar lake para GCS bucket",
        "Substituir FastAPI por Cloud Run",
        "Airflow gerenciado - Cloud Composer",
        "",
        "# Governanca",
        "OAuth no Airflow e MLflow",
        "Identity-Aware Proxy (IAP) no GCP",
    ])

    # Final
    slide_final(prs)

    prs.save(OUT_PPTX)
    print(f"\nSalvo: {OUT_PPTX}")
    print(f"{len(prs.slides)} slides com layouts FGV corretos.")


if __name__ == "__main__":
    main()
omocao automatizada",
        "Vector search em pgvector com indice HNSW cosine",
        "FastAPI servindo /predict (Registry) e /search (vetorial)",
        "Triplo versionamento: Git + lake parquet + MLflow Registry",
        "CI completo (ruff + pytest + docker build) via GitHub Actions",
        "Manual de deploy em VM GCP com script de provisionamento",
    ])
    slide_bullets(prs, "Proximos passos (evolucao)", [
        "# Curto prazo (laboratorio)",
        "Monitoracao de drift com Evidently AI (KS, PSI, Wasserstein)",
        "Datasheet for Datasets + Model Card por modelo registrado",
        "Testes E2E via docker compose + pytest-compose",
        "Avaliacao de fairness (Disparate Impact, Equal Opportunity)",
        "",
        "# Medio prazo (cloud)",
        "Migrar Postgres para Cloud SQL",
        "Migrar lake para GCS bucket",
        "Substituir FastAPI por Cloud Run",
        "Airflow gerenciado - Cloud Composer",
        "",
        "# Governanca",
        "OAuth no Airflow e MLflow",
        "Identity-Aware Proxy (IAP) no GCP",
    ])

    # Final
    slide_final(prs)

    prs.save(OUT_PPTX)
    print(f"\nSalvo: {OUT_PPTX}")
    print(f"{len(prs.slides)} slides com layouts FGV corretos.")


if __name__ == "__main__":
    main()
y (IAP) no GCP",
    ])

    # Final
    slide_final(prs)

    prs.save(OUT_PPTX)
    print(f"\nSalvo: {OUT_PPTX}")
    print(f"{len(prs.slides)} slides com layouts FGV corretos.")


if __name__ == "__main__":
    main()
