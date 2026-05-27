"""Apresentacao PPTX v2 - usando template real FGV + diagramas Graphviz/Mermaid.

- Carrega o template MLOps_E1_E2_DevOps_MLOps.pptx como base (herda design FGV).
- Renderiza diagramas como PNG via Graphviz (dot).
- Mantem .mmd source no docs/mermaid/ para refer.
- Adiciona NOVA secao 'Visao da Solucao' antes de problema com:
    - Entregaveis ao cliente final
    - Responsabilidades Airflow (orquestracao)
    - Responsabilidades MLflow (versionamento dados + modelos)
"""
from __future__ import annotations

import copy
import shutil
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# -----------------------------------------------------------------------------
TEMPLATE = Path("/sessions/busy-hopeful-clarke/mnt/uploads/MLOps_E1_ E2_DevOps_MLOps.pptx")
PROJECT_DIR = Path("/sessions/busy-hopeful-clarke/mnt/FGV/mlops-pratica")
DOCS_DIR = PROJECT_DIR / "docs"
IMG_DIR = DOCS_DIR / "img"
MMD_DIR = DOCS_DIR / "mermaid"
DOT_DIR = DOCS_DIR / "dot"
OUT_PPTX = PROJECT_DIR / "Apresentacao_MLOps_Pratica.pptx"

for d in [IMG_DIR, MMD_DIR, DOT_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# Cores FGV
AZUL = RGBColor(0x1F, 0x4E, 0x79)
AZUL_CLARO = RGBColor(0xD5, 0xE8, 0xF0)
VERMELHO = RGBColor(0xC0, 0x00, 0x00)
CINZA = RGBColor(0x59, 0x59, 0x59)
CINZA_ESC = RGBColor(0x3F, 0x3F, 0x3F)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

# Cores graphviz (strings hex)
HX = {
    "azul": "#1F4E79",
    "azul_claro": "#D5E8F0",
    "verde": "#5A8F4B",
    "vermelho": "#C00000",
    "cinza": "#595959",
    "laranja": "#ED7D31",
    "branco": "#FFFFFF",
    "amarelo_claro": "#FFF2CC",
}


# =============================================================================
# DIAGRAMAS - Mermaid source + DOT equivalente + PNG render via Graphviz
# =============================================================================
MERMAID_SRC = {
    "01_arquitetura": """flowchart LR
  API[HackerNews API] --> AF[Airflow]
  AF --> MIN[(MinIO data lake)]
  MIN --> MLF[MLflow Tracking + Registry]
  MIN --> M1[Modelo Preditivo]
  MIN --> M2[Modelo Embeddings]
  M1 --> MLF
  M2 --> MLF
  M2 --> PG[(Postgres pgvector)]
  MLF --> SRV[FastAPI /predict /search]
  PG --> SRV
""",
    "02_eng_dados": """flowchart LR
  A[extract HN topstories] --> B[raw parquet MinIO]
  B --> C[curate dedup]
  C --> D1[features tabular]
  C --> D2[features text]
  D1 --> E1[treino preditivo]
  D2 --> E2[treino embeddings]
""",
    "03_mlops": """flowchart LR
  L[Load features] --> T[Train RF + XGBoost]
  T --> E[Evaluate AUC/F1/PR-AUC]
  E --> M[MLflow Tracking]
  M --> R[Registry Staging]
  R --> P[Production]
  P --> API[FastAPI /predict]
""",
    "04_devops": """flowchart LR
  F[feature/*] --> D[develop]
  D --> R[release/*]
  R --> MA[main]
  MA --> T[tag v*]
  F -.-> CI1[CI ruff+pytest]
  R -.-> CI2[CI docker build]
  MA -.-> CD1[CD staging]
  T -.-> CD2[CD producao]
""",
    "05_dados": """flowchart LR
  API[HN API firebaseio] --> RAW[raw parquet]
  RAW --> CUR[curated parquet]
  CUR --> FT[features tabular]
  CUR --> FX[features text]
  FT --> PG[(Postgres pgvector)]
  FX --> PG
""",
    "06_gcp": """flowchart LR
  INET[Internet] --> FW[Firewall tag mlops]
  FW --> VM[Compute Engine e2-standard-4]
  VM --> DISK[(Persistent SSD 200GB)]
  VM -.-> CSQL[(Cloud SQL futuro)]
  VM -.-> GCS[(GCS bucket futuro)]
""",
    "07_solucao_geral": """flowchart TB
  subgraph CLIENT["O que o Cliente Final recebe"]
    A1[API REST FastAPI]
    A2[Endpoint /predict]
    A3[Endpoint /search semantico]
    A4[UIs operacionais Airflow MLflow MinIO]
  end
  subgraph PIPE["Backend MLOps"]
    B1[Airflow 3 DAGs]
    B2[MLflow Tracking+Registry]
    B3[MinIO data lake]
    B4[Postgres pgvector]
  end
  B1 --> B3
  B1 --> B2
  B2 --> B3
  B3 --> A2
  B4 --> A3
  B2 --> A2
""",
    "08_airflow_mlflow": """flowchart LR
  subgraph AF[Airflow ORQUESTRA]
    AF1[Quando rodar - schedule]
    AF2[Em que ordem - dependencias]
    AF3[Em caso de erro - retry/alert]
    AF4[Idempotencia das tasks]
  end
  subgraph ML[MLflow VERSIONA]
    ML1[Codigo do experimento]
    ML2[Hiperparametros e metricas]
    ML3[Artifacts do modelo]
    ML4[Stages do Registry]
  end
  subgraph LK[MinIO VERSIONA DADOS]
    LK1[raw por hora]
    LK2[curated dedup]
    LK3[features versionadas]
  end
  AF --> ML
  AF --> LK
""",
}


def write_mermaid_files():
    for name, src in MERMAID_SRC.items():
        (MMD_DIR / f"{name}.mmd").write_text(src, encoding="utf-8")


# -----------------------------------------------------------------------------
# Conversao para DOT (Graphviz) - estilo flowchart limpo similar a Mermaid
# -----------------------------------------------------------------------------
def dot_header(rankdir="LR"):
    return f"""digraph G {{
    rankdir={rankdir};
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=12,
          fontcolor="white", margin="0.20,0.12"];
    edge [color="{HX['cinza']}", penwidth=1.6, arrowsize=0.9];
    splines=spline;
    nodesep=0.5;
    ranksep=0.7;
"""


def dot_node(nid, label, fill="#1F4E79", shape="box", fontcolor="white"):
    label = label.replace("\n", "\\n")
    return (f'    {nid} [label="{label}", fillcolor="{fill}", shape={shape}, '
            f'fontcolor="{fontcolor}"];\n')


def dot_edge(a, b, dashed=False, label=None):
    style = ', style="dashed"' if dashed else ""
    lab = f', label="{label}"' if label else ""
    return f'    {a} -> {b} [{style[2:] if style else ""}{lab}];\n'.replace("[, ", "[")


def build_dot_diagrams() -> dict:
    """Constroi diagramas DOT correspondentes aos Mermaid."""
    diags = {}

    # 01 arquitetura - LR
    d = dot_header("LR")
    d += dot_node("API", "HackerNews\\nAPI", fill=HX["verde"])
    d += dot_node("AF", "Airflow\\nweb + scheduler", fill=HX["azul"])
    d += dot_node("MIN", "MinIO\\ndata lake S3", fill=HX["azul"], shape="cylinder")
    d += dot_node("MLF", "MLflow\\nTracking + Registry", fill=HX["azul"])
    d += dot_node("M1", "ML Preditivo\\nsklearn + XGBoost", fill=HX["cinza"])
    d += dot_node("M2", "Embeddings\\nMiniLM-L6 (CPU)", fill=HX["cinza"])
    d += dot_node("PG", "Postgres + pgvector", fill=HX["laranja"], shape="cylinder")
    d += dot_node("SRV", "FastAPI\\n/predict /search", fill=HX["azul"])
    for a, b in [("API", "AF"), ("AF", "MIN"), ("MIN", "MLF"),
                 ("MIN", "M1"), ("MIN", "M2"), ("M1", "MLF"),
                 ("M2", "MLF"), ("M2", "PG"), ("MLF", "SRV"), ("PG", "SRV")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["01_arquitetura"] = d

    # 02 engenharia de dados - LR
    d = dot_header("LR")
    d += dot_node("E", "1. extract\\nHN topstories", fill=HX["verde"])
    d += dot_node("R", "2. raw\\nparquet", fill=HX["azul"])
    d += dot_node("C", "3. curate\\ndedup por id", fill=HX["azul"])
    d += dot_node("FT", "4a. features\\ntabular", fill=HX["laranja"])
    d += dot_node("FX", "4b. features\\ntext", fill=HX["laranja"])
    d += dot_node("TP", "treino preditivo", fill=HX["cinza"])
    d += dot_node("TE", "treino embeddings", fill=HX["cinza"])
    for a, b in [("E", "R"), ("R", "C"), ("C", "FT"), ("C", "FX"),
                 ("FT", "TP"), ("FX", "TE")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["02_eng_dados"] = d

    # 03 MLOps - LR
    d = dot_header("LR")
    d += dot_node("L", "Load features\\n(MinIO)", fill=HX["azul"])
    d += dot_node("T", "Train\\nRF + XGBoost", fill=HX["azul"])
    d += dot_node("EV", "Evaluate\\nAUC/F1/PR-AUC", fill=HX["cinza"])
    d += dot_node("M", "MLflow Tracking\\nparams + metrics\\nartifacts + signature", fill=HX["laranja"])
    d += dot_node("R", "Registry\\nNone -> Staging", fill=HX["azul"])
    d += dot_node("P", "Promote\\n-> Production", fill=HX["verde"])
    d += dot_node("API", "FastAPI\\n/predict /search", fill=HX["azul"])
    for a, b in [("L", "T"), ("T", "EV"), ("EV", "M"),
                 ("M", "R"), ("R", "P"), ("P", "API")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["03_mlops"] = d

    # 04 DevOps - LR
    d = dot_header("LR")
    d += dot_node("F", "feature/*", fill=HX["verde"])
    d += dot_node("D", "develop", fill=HX["laranja"])
    d += dot_node("R", "release/*", fill=HX["cinza"])
    d += dot_node("MA", "main", fill=HX["azul"])
    d += dot_node("T", "tag v*", fill=HX["azul"])
    d += dot_node("CI1", "CI\\nruff + pytest", fill=HX["verde"], fontcolor="white")
    d += dot_node("CI2", "CI\\ndocker build", fill=HX["cinza"])
    d += dot_node("CD1", "CD staging", fill=HX["azul"])
    d += dot_node("CD2", "CD producao", fill=HX["laranja"])
    for a, b in [("F", "D"), ("D", "R"), ("R", "MA"), ("MA", "T")]:
        d += f"    {a} -> {b};\n"
    for a, b in [("F", "CI1"), ("R", "CI2"), ("MA", "CD1"), ("T", "CD2")]:
        d += f'    {a} -> {b} [style="dashed"];\n'
    d += "}"
    diags["04_devops"] = d

    # 05 dados - LR
    d = dot_header("LR")
    d += dot_node("API", "HN API\\nfirebaseio", fill=HX["verde"])
    d += dot_node("RAW", "raw\\nparquet MinIO\\ndt=YYYY-MM-DD/hr=HH", fill=HX["azul"])
    d += dot_node("CUR", "curated\\ndedup por id", fill=HX["azul"])
    d += dot_node("FT", "features tabular", fill=HX["laranja"])
    d += dot_node("FX", "features text", fill=HX["laranja"])
    d += dot_node("PG", "Postgres + pgvector\\nairflow / mlflow / app", fill=HX["laranja"], shape="cylinder")
    for a, b in [("API", "RAW"), ("RAW", "CUR"), ("CUR", "FT"),
                 ("CUR", "FX"), ("FT", "PG"), ("FX", "PG")]:
        d += f"    {a} -> {b};\n"
    d += "}"
    diags["05_dados"] = d

    # 06 GCP - LR
    d = dot_header("LR")
    d += dot_node("INET", "Internet", fill=HX["cinza"])
    d += dot_node("FW", "Firewall\\ntag mlops\\n8080/5000/9001/8000", fill=HX["laranja"])
    d += dot_node("VM", "Compute Engine\\ne2-standard-4\\nUbuntu 22.04", fill=HX["azul"])
    d += dot_node("DISK", "Persistent SSD\\n200GB", fill=HX["cinza"], shape="cylinder")
    d += dot_node("CSQL", "Cloud SQL\\n(evolucao)", fill=HX["verde"], shape="cylinder")
    d += dot_node("GCS", "GCS bucket\\n(evolucao)", fill=HX["verde"], shape="cylinder")
    for a, b in [("INET", "FW"), ("FW", "VM"), ("VM", "DISK")]:
        d += f"    {a} -> {b};\n"
    for a, b in [("VM", "CSQL"), ("VM", "GCS")]:
        d += f'    {a} -> {b} [style="dashed"];\n'
    d += "}"
    diags["06_gcp"] = d

    # 07 solucao geral - TB com clusters
    d = """digraph G {
    rankdir=TB;
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=12,
          fontcolor="white", margin="0.20,0.12"];
    edge [color="#595959", penwidth=1.6, arrowsize=0.9];
    nodesep=0.4; ranksep=0.6;

    subgraph cluster_pipe {
        label="Backend MLOps (interno - tudo Docker)";
        labeljust=l; labelloc=t;
        style="rounded,filled"; fillcolor="#EAF3F8"; color="#1F4E79";
        fontcolor="#1F4E79"; fontsize=14;
        B1 [label="Airflow\\n3 DAGs", fillcolor="#1F4E79"];
        B2 [label="MLflow\\nTracking + Registry", fillcolor="#1F4E79"];
        B3 [label="MinIO\\ndata lake", fillcolor="#1F4E79"];
        B4 [label="Postgres\\npgvector", fillcolor="#ED7D31"];
    }

    subgraph cluster_client {
        label="O que o Cliente Final usa";
        labeljust=l; labelloc=t;
        style="rounded,filled"; fillcolor="#FFF2CC"; color="#C28A00";
        fontcolor="#C28A00"; fontsize=14;
        A1 [label="FastAPI\\n/predict", fillcolor="#5A8F4B"];
        A2 [label="FastAPI\\n/search semantico", fillcolor="#5A8F4B"];
        A3 [label="UIs operacionais\\nAirflow / MLflow / MinIO", fillcolor="#5A8F4B"];
    }

    B1 -> B3; B1 -> B2; B2 -> B3; B3 -> A1; B4 -> A2; B2 -> A1; B3 -> A3; B2 -> A3;
}
"""
    diags["07_solucao_geral"] = d

    # 08 Airflow vs MLflow vs MinIO
    d = """digraph G {
    rankdir=LR;
    bgcolor="white";
    node [shape=box, style="filled,rounded", fontname="Helvetica", fontsize=11,
          fontcolor="white", margin="0.18,0.10"];
    edge [color="#595959", penwidth=1.5];
    nodesep=0.3; ranksep=0.8;

    subgraph cluster_af {
        label="Airflow ORQUESTRA"; labeljust=l;
        style="rounded,filled"; fillcolor="#EAF3F8"; color="#1F4E79";
        fontcolor="#1F4E79"; fontsize=14;
        AF1 [label="Quando rodar\\n(schedule cron)", fillcolor="#1F4E79"];
        AF2 [label="Em que ordem\\n(dependencias)", fillcolor="#1F4E79"];
        AF3 [label="Retry e alerta\\nem caso de erro", fillcolor="#1F4E79"];
        AF4 [label="Idempotencia\\ndas tasks", fillcolor="#1F4E79"];
    }
    subgraph cluster_ml {
        label="MLflow VERSIONA MODELOS"; labeljust=l;
        style="rounded,filled"; fillcolor="#FFF2CC"; color="#C28A00";
        fontcolor="#C28A00"; fontsize=14;
        ML1 [label="Codigo do\\nexperimento (run)", fillcolor="#ED7D31"];
        ML2 [label="Params + metrics\\n(track)", fillcolor="#ED7D31"];
        ML3 [label="Artifacts\\n(pipeline pickle)", fillcolor="#ED7D31"];
        ML4 [label="Stages do Registry\\n(None/Staging/Prod)", fillcolor="#ED7D31"];
    }
    subgraph cluster_lk {
        label="MinIO VERSIONA DADOS"; labeljust=l;
        style="rounded,filled"; fillcolor="#E2F0D9"; color="#5A8F4B";
        fontcolor="#5A8F4B"; fontsize=14;
        LK1 [label="raw por hora\\n(append imutavel)", fillcolor="#5A8F4B"];
        LK2 [label="curated\\n(snapshot dedup)", fillcolor="#5A8F4B"];
        LK3 [label="features\\nversionadas", fillcolor="#5A8F4B"];
    }
    AF1 -> ML1 [style=invis];
    AF1 -> LK1 [style=invis];
}
"""
    diags["08_airflow_mlflow"] = d

    return diags


def render_dot_to_png(dots: dict) -> dict:
    """Renderiza cada DOT como PNG com dot -Tpng."""
    images = {}
    for name, content in dots.items():
        dot_path = DOT_DIR / f"{name}.dot"
        png_path = IMG_DIR / f"{name}.png"
        dot_path.write_text(content, encoding="utf-8")
        result = subprocess.run(
            ["dot", "-Tpng", "-Gdpi=180", str(dot_path), "-o", str(png_path)],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(f"dot falhou para {name}: {result.stderr}")
        images[name] = str(png_path)
        print(f"  -> {png_path}")
    return images


# =============================================================================
# PPTX - usando o template como base
# =============================================================================
def clone_template():
    """Carrega o template e remove todos os slides existentes."""
    prs = Presentation(str(TEMPLATE))
    # Remove slides existentes (mas mantem masters/layouts)
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for sldId in slides_to_remove:
        rId = sldId.attrib[qn("r:id")]
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)
    print(f"Template carregado. Slides apagados. Layouts disponiveis: {len(prs.slide_layouts)}")
    return prs


def get_layout(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise ValueError(f"Layout '{name}' nao encontrado")


def _set_run(run, text, size=14, bold=False, color=None, name="Calibri", italic=False):
    run.text = text
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def set_placeholder_text(slide, idx, text, size=14, bold=False, color=None, italic=False, align=None):
    """Preenche um placeholder pelo idx. Aceita string ou lista de paragrafos."""
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        _set_run(r, line, size=size, bold=bold, color=color, italic=italic)


def add_bullets_to_placeholder(slide, idx, lines):
    """Bullets formatados. Linhas que comecam com '# ' viram subtitulos azuis,
    linhas vazias = espacador, demais = bullets."""
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        if line.startswith("# "):
            _set_run(r, line[2:], size=18, bold=True, color=AZUL)
            p.space_after = Pt(4)
        elif line == "":
            _set_run(r, " ", size=8)
            p.space_after = Pt(2)
        else:
            _set_run(r, "• " + line, size=14, color=CINZA_ESC)
            p.space_after = Pt(4)


def add_picture_centered(slide, img_path, left=None, top=None,
                         width=None, height=None):
    pic = slide.shapes.add_picture(img_path, left or 0, top or 0,
                                   width=width, height=height)
    return pic


# -----------------------------------------------------------------------------
# Slides
# -----------------------------------------------------------------------------
SW = Inches(20)  # template slide width
SH = Inches(11.25)


def slide_capa(prs):
    s = prs.slides.add_slide(get_layout(prs, "Slide de abertura"))
    # Title placeholder idx=0, body idx=10
    try:
        set_placeholder_text(s, 0, "Pratica MLOps end-to-end", size=44, bold=True, color=AZUL)
        set_placeholder_text(s, 10,
            ["HackerNews + Airflow + MLflow + pgvector + FastAPI",
             "100% Docker | Deploy local e em VM GCP",
             "Prof. Andre Insardi — Maio/2026"],
            size=18, color=CINZA)
    except Exception:
        pass
    return s


def slide_titulo_conteudo(prs, titulo, bullets, subtitulo=None):
    s = prs.slides.add_slide(get_layout(prs, "1_Title and Content"))
    set_placeholder_text(s, 0, titulo, size=28, bold=True, color=AZUL)
    add_bullets_to_placeholder(s, 1, bullets)
    return s


def slide_imagem_full(prs, titulo, img_path, legenda=None):
    s = prs.slides.add_slide(get_layout(prs, "1_Title and Content"))
    set_placeholder_text(s, 0, titulo, size=28, bold=True, color=AZUL)
    # Substitui o body placeholder por imagem grande
    body = s.placeholders[1]
    # Remove o placeholder e adiciona imagem
    sp = body._element
    sp.getparent().remove(sp)
    # Imagem centralizada
    s.shapes.add_picture(img_path, Inches(2.5), Inches(2.0), width=Inches(15.0))
    if legenda:
        tb = s.shapes.add_textbox(Inches(2.5), Inches(9.8), Inches(15.0), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        _set_run(r, legenda, size=12, italic=True, color=CINZA)
    return s


def slide_divisor(prs, numero, titulo, descricao):
    """Slide divisor (capa de bloco)."""
    s = prs.slides.add_slide(get_layout(prs, "Blank"))
    # Fundo claro
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xEA, 0xF3, 0xF8)
    # Numero grande
    tb = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(4.0), Inches(3.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, numero, size=200, bold=True, color=AZUL)
    # Titulo
    tb = s.shapes.add_textbox(Inches(6.0), Inches(3.8), Inches(13.0), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, titulo, size=44, bold=True, color=AZUL)
    # Linha vermelha
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(5.2),
                              Inches(3.0), Inches(0.08))
    line.fill.solid(); line.fill.fore_color.rgb = VERMELHO
    line.line.fill.background()
    # Descricao
    tb = s.shapes.add_textbox(Inches(6.0), Inches(5.5), Inches(13.0), Inches(3.0))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, descricao, size=20, color=CINZA_ESC)
    return s


def slide_final(prs):
    s = prs.slides.add_slide(get_layout(prs, "1_slideFinal"))
    try:
        set_placeholder_text(s, 15,
            ["Obrigado!",
             "ext.andre.insardi@prof.fgv.edu.br",
             "FGV — MBA IA & Analytics — MLOps — Maio/2026"],
            size=22, bold=True, color=BRANCO)
    except Exception:
        pass
    return s


# =============================================================================
def main():
    print("==> Salvando .mmd sources")
    write_mermaid_files()
    print("==> Construindo DOTs e renderizando PNGs com Graphviz")
    dots = build_dot_diagrams()
    imgs = render_dot_to_png(dots)

    print("==> Carregando template e zerando slides")
    prs = clone_template()

    print("==> Montando slides")

    # 1 - Capa
    slide_capa(prs)

    # 2 - Agenda
    slide_titulo_conteudo(
        prs,
        "Agenda da Pratica Integradora",
        [
            "# Roteiro (9 blocos)",
            "0. Visao da solucao - o que vamos construir e entregar",
            "1. Problema e caso de uso",
            "2. Base de dados",
            "3. Arquitetura geral",
            "4. Engenharia de Dados (Airflow)",
            "5. Fluxo MLOps (MLflow)",
            "6. Fluxo DevOps (Gitflow)",
            "7. Deploy em VM GCP",
            "8. Conclusoes e proximos passos",
        ],
    )

    # === BLOCO 0 NOVO - Visao da Solucao =====================================
    slide_divisor(
        prs, "0", "Visao da Solucao",
        "Antes de mergulhar nos detalhes: o que sera entregue ao cliente final, "
        "o que o Airflow controla, o que o MLflow versiona e como os dados sao versionados no lake."
    )

    slide_titulo_conteudo(
        prs,
        "O que vamos construir (resumo executivo)",
        [
            "# Visao em uma frase",
            "Uma plataforma MLOps end-to-end rodando em Docker (laptop ou VM GCP) que ingere dados publicos do HackerNews, treina dois pipes de ML (preditivo + embeddings) e expoe servicos via API REST.",
            "",
            "# Componentes principais (7 containers)",
            "Postgres + pgvector (metadata + vector store).",
            "MinIO (data lake S3-compatible).",
            "MLflow (Tracking + Model Registry).",
            "Airflow webserver + scheduler (orquestracao das 3 DAGs).",
            "FastAPI (serving /predict e /search semantico).",
            "",
            "# Tres pipelines orquestrados",
            "Ingestao horaria de stories do HN.",
            "Treino diario do classificador 'vai bombar' (sklearn + XGBoost).",
            "Treino diario de embeddings dos titulos (MiniLM-L6) com indexacao em pgvector.",
        ],
    )

    slide_imagem_full(
        prs,
        "Visao end-to-end - Backend MLOps + Cliente Final",
        imgs["07_solucao_geral"],
        legenda="O backend (azul) fica oculto ao cliente; o que ele consome sao as APIs e UIs (verde).",
    )

    slide_titulo_conteudo(
        prs,
        "O que o Cliente Final recebe",
        [
            "# Servicos de uso direto (API REST)",
            "POST /predict - dado um post novo (titulo, hora, autor, URL), devolve probabilidade de 'vai bombar' (score >= 100).",
            "POST /search - busca semantica de posts similares por similaridade cosseno em pgvector (top-K configuravel).",
            "GET  /health - liveness probe (model_loaded + tracking_uri).",
            "",
            "# Interfaces operacionais (web UIs)",
            "Airflow UI (8080) - status das DAGs, logs, retry manual.",
            "MLflow UI (5000) - experimentos, runs, model registry e versoes.",
            "MinIO Console (9001) - exploracao do data lake (raw / curated / features).",
            "Swagger UI (8000/docs) - documentacao interativa da API.",
            "",
            "# SLAs implicitos (didaticos)",
            "Latencia /predict: < 100ms (modelo carregado em memoria).",
            "Latencia /search: < 200ms para top-10 em ate ~1M vetores no pgvector.",
            "Freshness: dados ingeridos a cada hora; modelo retreinado diariamente.",
            "Disponibilidade: depende do host (laptop = best-effort; VM GCP = ~99%).",
        ],
    )

    slide_imagem_full(
        prs,
        "Quem controla o que - Airflow vs MLflow vs MinIO",
        imgs["08_airflow_mlflow"],
        legenda="Airflow orquestra o quando/ordem. MLflow versiona modelos. MinIO versiona dados.",
    )

    slide_titulo_conteudo(
        prs,
        "O que o Airflow controla (orquestracao)",
        [
            "# Responsabilidades exclusivas do Airflow",
            "QUANDO cada pipeline roda - via schedule (cron expression).",
            "EM QUE ORDEM as tasks executam - via dependencias declarativas no DAG.",
            "EM CASO DE ERRO - retry com backoff exponencial e alertas configuraveis.",
            "IDEMPOTENCIA - cada task pode ser re-executada sem corromper o estado.",
            "OBSERVABILIDADE - logs, tempo de execucao e historico de runs por task.",
            "",
            "# Schedule das 3 DAGs no projeto",
            "pipeline_ingestao @hourly - extract + raw + curated + features.",
            "pipeline_treino_preditivo 02:00 daily - load + train RF/XGBoost + register.",
            "pipeline_treino_embeddings 03:00 daily - load + encode MiniLM + upsert pgvector.",
            "",
            "# O que o Airflow NAO faz",
            "Nao versiona codigo (isso e papel do Git).",
            "Nao versiona modelos (isso e papel do MLflow Registry).",
            "Nao armazena features (isso e papel do MinIO/lake).",
            "Nao serve predicoes (isso e papel do FastAPI).",
        ],
    )

    slide_titulo_conteudo(
        prs,
        "O que o MLflow versiona (modelos)",
        [
            "# Tracking (durante o treino)",
            "params - hiperparametros do estimator (n_estimators, max_depth, learning_rate, ...).",
            "metrics - accuracy, F1, ROC-AUC, PR-AUC em test split estratificado.",
            "artifacts - pipeline sklearn completo (preprocessor + estimator).",
            "signature - schema de input/output inferido automaticamente.",
            "tags - modelo, dataset_version, autor, ambiente.",
            "",
            "# Model Registry (gestao de versoes)",
            "Cada run que registra um modelo cria uma nova VERSAO em hn_classifier.",
            "Stages: None -> Staging -> Production -> Archived.",
            "Promocao programatica (DAG) ou via UI com aprovacao humana.",
            "FastAPI carrega o modelo por stage: models:/hn_classifier/Production.",
            "",
            "# Dois experimentos separados",
            "Experimento 'hn_classifier' - runs do pipeline preditivo (RF + XGBoost).",
            "Experimento 'hn_embeddings' - runs do encoder MiniLM (versao, dim, amostras).",
            "",
            "# O que o MLflow NAO faz",
            "Nao orquestra DAGs (Airflow).",
            "Nao armazena dados brutos (MinIO/lake).",
            "Nao executa o treino - apenas RECEBE logs durante o treino.",
        ],
    )

    slide_titulo_conteudo(
        prs,
        "O que o MinIO versiona (dados)",
        [
            "# Tres camadas do lake",
            "raw - cada execucao da DAG horaria grava uma particao imutavel (dt=YYYY-MM-DD/hr=HH).",
            "curated - snapshot estavel deduplicado por id (mantem maior score visto).",
            "features - dataset materializado para treino (tabular e text).",
            "",
            "# Como a versao acontece",
            "raw: nunca sobrescreve - cada hora gera um arquivo NOVO (auditavel).",
            "curated: append-merge com dedup; e o dado canonico do dia.",
            "features: cada DAG de treino le a versao mais recente.",
            "",
            "# Caminhos completos",
            "s3://raw/hn/dt=2026-05-26/hr=14/items.parquet",
            "s3://curated/hn/stories.parquet",
            "s3://features/tabular/train.parquet  /  s3://features/text/titles.parquet",
            "",
            "# Evolucao para DVC ou Delta Lake (futuro)",
            "Para versionamento explicito com hash + lineage, evoluir features para DVC ou Delta Lake.",
            "No estado atual, o particionamento por data ja garante reprodutibilidade temporal.",
        ],
    )

    slide_titulo_conteudo(
        prs,
        "Como Airflow, MLflow e MinIO conversam em runtime",
        [
            "# Fluxo de uma execucao de treino preditivo",
            "1. Scheduler Airflow detecta horario (02:00) e dispara DAG.",
            "2. Task 'load_features' le features/tabular/train.parquet do MinIO via s3fs.",
            "3. Task 'train' usa MLflow client: mlflow.set_experiment('hn_classifier') + start_run().",
            "4. Durante o treino, MLflow loga params/metrics no Postgres (backend store).",
            "5. mlflow.sklearn.log_model() envia o pipeline pickle para MinIO via S3 API (artifact store).",
            "6. mlflow.register_model() cria a versao N do modelo hn_classifier.",
            "7. Task 'promote' transiciona via REST API: stage=Staging.",
            "",
            "# Fluxo de uma predicao em runtime",
            "1. Cliente POST /predict no FastAPI.",
            "2. FastAPI consulta MLflow Registry: models:/hn_classifier/Production.",
            "3. MLflow lookup retorna URI s3://mlflow-artifacts/.../model.",
            "4. FastAPI baixa pickle do MinIO, carrega em memoria (cache no startup).",
            "5. Aplica pipeline e devolve probabilidade. < 100ms tipico.",
            "",
            "# Disclaimer didatico",
            "Em producao: artifact store separado (S3 nativo + CDN), Registry com gates de aprovacao, monitoracao de drift, retraining automatico por threshold de degradacao.",
        ],
    )

    # === BLOCO 1 - Problema (mais detalhado) =================================
    slide_divisor(
        prs, "1", "Problema",
        "Por que precisamos de uma pratica integradora, qual e a dor pedagogica "
        "que estamos resolvendo e qual e o caso de uso real escolhido.",
    )
    slide_titulo_conteudo(
        prs,
        "Contexto pedagogico - por que essa pratica existe",
        [
            "# Diagnostico nos encontros E1-E4",
            "E1 ensinou fundamentos (DevOps -> MLOps, dívida tecnica, ciclo de vida).",
            "E2 mostrou experimentacao (autolog, Registry, HPO, Model Card).",
            "E3 trabalhou producao (deploy patterns, FastAPI/Docker, CI/CD, drift).",
            "E4 cobriu governanca (LGPD, AI Act, LLMOps, RAG).",
            "",
            "# Gap identificado",
            "Cada encontro discutiu blocos isolados. Falta INTEGRAR.",
            "Alunos precisam VER o ciclo completo executando: do trigger ao serving.",
            "Querem entender RESPONSABILIDADES (quem orquestra, quem versiona, quem serve).",
            "",
            "# O que esta pratica entrega ao aluno",
            "Repositorio rodavel no laptop com 'make up'.",
            "Mesma stack subindo em VM GCP (mesmo compose + override de volumes).",
            "Codigo Python modular e documentado (PEP-8 + ruff + pytest).",
            "Apresentacao narrando cada bloco com diagramas Mermaid.",
        ],
    )
    slide_titulo_conteudo(
        prs,
        "Caso de uso - HackerNews 'vai bombar'",
        [
            "# Definicao do target",
            "vai_bombar = (score >= 100) - threshold didatico que gera ~10-20% positivos.",
            "Alternativas avaliadas: regressao log(score), classificacao multiclasse por tipo.",
            "Binario foi escolhido por simplicidade pedagogica e metricas intuitivas (AUC, F1).",
            "",
            "# Por que HackerNews?",
            "API publica sem autenticacao (zero atrito de credenciais).",
            "Dado heterogeneo: tabular (score, comments, hora) + texto (titulo, URL).",
            "Volume controlavel (top 100 stories por extracao).",
            "Caso analogo ao real: priorizar feed editorial, alerta de conteudo viral.",
            "",
            "# Features que vamos usar",
            "Tabular: title_len, n_words_title, hour, weekday, has_url, has_question, domain, by_author.",
            "Vetorial: embedding 384-d do titulo via sentence-transformers/all-MiniLM-L6-v2.",
            "",
            "# Aplicacoes empresariais analogas",
            "Triagem de tickets de suporte; priorizacao de leads; detecccao de churn risk; deduplicacao de topicos em mídia.",
        ],
    )

    # === BLOCO 2 - Base de dados =============================================
    slide_divisor(prs, "2", "Base de dados",
                  "Fonte, camadas do lake e modelagem - HackerNews API + MinIO + Postgres/pgvector.")
    slide_titulo_conteudo(
        prs, "Fonte e camadas de dados",
        [
            "# Fonte",
            "HackerNews Firebase API (publica, sem auth, sem rate limit pratico).",
            "Endpoints: /v0/topstories.json (lista IDs) e /v0/item/{id}.json (detalhe).",
            "",
            "# Lake (MinIO, formato parquet)",
            "raw/hn/dt=YYYY-MM-DD/hr=HH/items.parquet (particionamento Hive-style).",
            "curated/hn/stories.parquet (dedup por id, snapshot do score maximo).",
            "features/tabular/train.parquet (input do pipe preditivo).",
            "features/text/titles.parquet (input do pipe de embeddings).",
            "",
            "# Postgres (3 databases na mesma instancia)",
            "airflow - metadata do scheduler.",
            "mlflow - backend store do Tracking Server (runs + Registry).",
            "app - tabela embeddings(id, title, vector(384), score, ts) + indice HNSW vector_cosine_ops.",
        ],
    )
    slide_imagem_full(prs, "Diagrama da base de dados", imgs["05_dados"],
                      legenda="Camadas raw, curated e features no MinIO + Postgres com pgvector.")

    # === BLOCO 3 - Arquitetura ===============================================
    slide_divisor(prs, "3", "Arquitetura geral",
                  "Sete containers Docker em uma rede unica, sem dependencias de cloud.")
    slide_titulo_conteudo(
        prs, "Stack confirmada (7 containers)",
        [
            "# Camadas e tecnologias",
            "Orquestracao - Apache Airflow 2.9 (LocalExecutor).",
            "Tracking + Registry - MLflow 2.16 (backend Postgres, artifact MinIO).",
            "Data Lake - MinIO (S3-compatible).",
            "Metadata + Vector - Postgres 16 + pgvector 0.7 (HNSW cosine).",
            "Embeddings - sentence-transformers/all-MiniLM-L6-v2 (offline, CPU, dim=384).",
            "ML Preditivo - scikit-learn 1.5 + XGBoost 2.1.",
            "Serving - FastAPI 0.115 + Uvicorn.",
            "CI/CD - GitHub Actions (ruff + pytest + docker build).",
        ],
    )
    slide_imagem_full(prs, "Arquitetura geral - 7 containers em rede unica", imgs["01_arquitetura"],
                      legenda="Postgres + MinIO + MLflow + Airflow (web+sched) + FastAPI.")

    # === BLOCO 4 - Engenharia de Dados =======================================
    slide_divisor(prs, "4", "Engenharia de Dados (Airflow)",
                  "DAGs, particionamento, idempotencia e features.")
    slide_imagem_full(prs, "DAG pipeline_ingestao (@hourly)", imgs["02_eng_dados"],
                      legenda="extract -> raw -> curated -> features (tabular + text).")
    slide_titulo_conteudo(
        prs, "Padroes de design das DAGs",
        [
            "# TaskFlow API",
            "Decorators @dag e @task (Python puro, sem operadores legados).",
            "XCom automatico via return + assinatura entre tasks.",
            "",
            "# Idempotencia e particionamento",
            "Cada execucao sobrescreve a particao da hora corrente.",
            "Particionamento Hive-style (dt=, hr=) habilita query eficiente futura.",
            "",
            "# Dedup na curated",
            "Concat (curated + raw novo) -> sort score desc -> drop_duplicates(id).",
            "",
            "# Resilience",
            "Cliente HTTP com tenacity (3 retries, backoff exponencial).",
            "ThreadPoolExecutor 16 workers para hidratar 100 items em ~5s.",
            "Retry da task (2x, 2min) cobre falhas transientes da API.",
        ],
    )

    # === BLOCO 5 - MLOps =====================================================
    slide_divisor(prs, "5", "Fluxo MLOps (MLflow)",
                  "Tracking, Registry, promocao para Production e serving via FastAPI.")
    slide_imagem_full(prs, "Fluxo MLOps - tracking, registry e serving", imgs["03_mlops"],
                      legenda="Dois experimentos -> Registry -> FastAPI /predict.")
    slide_titulo_conteudo(
        prs, "Triplo versionamento",
        [
            "# 1. Codigo - Git",
            "Branches main (estavel), develop (integracao), feature/* e release/*.",
            "Tags v* para releases (v0.1, v0.2, ...).",
            "",
            "# 2. Dados - Parquet particionado em MinIO",
            "raw/ guarda o estado bruto da hora (auditavel).",
            "curated/ e a versao estavel (snapshot).",
            "features/ versiona o dataset de treino (evolucao -> DVC ou Delta).",
            "",
            "# 3. Modelos - MLflow Tracking + Registry",
            "Run captura params, metrics, artifacts e signature.",
            "Versoes numericas + stages (None | Staging | Production | Archived).",
        ],
    )
    slide_titulo_conteudo(
        prs, "DAG treino preditivo (02:00 daily)",
        [
            "# Etapas",
            "load_features -> train_and_register_models -> promote.",
            "",
            "# Modelos comparados",
            "RandomForest (n_estimators=200, max_depth=12).",
            "XGBoost (n_estimators=300, max_depth=6, learning_rate=0.1).",
            "Vencedor escolhido por ROC-AUC.",
            "",
            "# Politica de promocao",
            "DAG promove ultima versao para Staging automaticamente.",
            "Producao real: gate humano + comparativo de metricas no Registry.",
        ],
    )
    slide_titulo_conteudo(
        prs, "DAG treino embeddings (03:00 daily)",
        [
            "# Etapas",
            "load_text -> encode -> index.",
            "",
            "# Encoder",
            "sentence-transformers/all-MiniLM-L6-v2 (offline, CPU, dim=384).",
            "Pre-baixado no Dockerfile (evita download em runtime).",
            "Normalizacao L2 dos vetores (cosine distance correto).",
            "",
            "# Indexacao",
            "Upsert em embeddings(id, title, vector, score, by_author, ts).",
            "Indice HNSW vector_cosine_ops (rapido para milhoes de vetores).",
        ],
    )

    # === BLOCO 6 - DevOps ====================================================
    slide_divisor(prs, "6", "Fluxo DevOps (Gitflow)",
                  "Branches, CI no GitHub Actions e CD por estagio.")
    slide_imagem_full(prs, "Gitflow + CI/CD", imgs["04_devops"],
                      legenda="GitHub Actions roda ruff + pytest em PRs e docker build em release.")
    slide_titulo_conteudo(
        prs, "Workflow ci.yml (GitHub Actions)",
        [
            "# Triggers",
            "push em main/develop e pull_request em main/develop.",
            "",
            "# Job 1 - lint-test",
            "Setup Python 3.11 + pip cache.",
            "ruff check src tests dags.",
            "pytest --cov=src/mlops_pratica --cov-report=term-missing.",
            "",
            "# Job 2 - docker-build (needs lint-test)",
            "docker buildx setup.",
            "Build sanity-check das 3 imagens: airflow, mlflow, fastapi.",
            "",
            "# Politica de release",
            "Tag vX.Y na main dispara workflow de release (pode habilitar push para GHCR).",
            "Gate humano antes de qualquer deploy compartilhado.",
        ],
    )

    # === BLOCO 7 - Deploy GCP ================================================
    slide_divisor(prs, "7", "Deploy em VM GCP",
                  "Compute Engine + Persistent Disk + Firewall via script gcloud.")
    slide_imagem_full(prs, "Arquitetura GCP (VM unica)", imgs["06_gcp"],
                      legenda="Firewall por tag mlops libera as 4 portas das UIs.")
    slide_titulo_conteudo(
        prs, "Como provisionar",
        [
            "# Pre-requisitos",
            "Projeto GCP com billing.",
            "gcloud autenticado e cota para 4 vCPU.",
            "",
            "# Comando unico",
            "export GCP_PROJECT=meu-projeto",
            "./infra/gcp/provision.sh",
            "",
            "# O script faz",
            "Cria disco SSD 200GB (pd-ssd).",
            "Cria VM ubuntu-22 com startup-script.",
            "Abre firewall por tag (ajuste SOURCE_IP_RANGE em prod!).",
            "",
            "# Startup-script automatiza",
            "Instala Docker Engine + Compose plugin.",
            "Formata e monta disco em /opt/mlops.",
            "Clona o repositorio e roda docker compose up -d.",
            "",
            "# Custo estimado",
            "~US$ 150/mes (VM 24/7 + 300GB SSD). Preemptible reduz ~70%.",
        ],
    )

    # === BLOCO 8 - Conclusoes ================================================
    slide_divisor(prs, "8", "Conclusoes",
                  "O que entregamos e proximos passos para evolucao.")
    slide_titulo_conteudo(
        prs, "O que entregamos",
        [
            "# Componentes implementados",
            "7 containers Docker orquestrados via docker-compose.",
            "3 DAGs Airflow cobrindo ingestao, treino preditivo e embeddings.",
            "MLflow Tracking + Registry com promocao automatizada.",
            "Vector search em pgvector com indice HNSW cosine.",
            "FastAPI servindo /predict (Registry) e /search (vetorial).",
            "Triplo versionamento: Git + lake parquet + MLflow Registry.",
            "CI completo (ruff + pytest + docker build) via GitHub Actions.",
            "Manual de deploy em VM GCP com script de provisionamento.",
        ],
    )
    slide_titulo_conteudo(
        prs, "Proximos passos (evolucao)",
        [
            "# Curto prazo (laboratorio)",
            "Monitoracao de drift com Evidently AI (KS, PSI, Wasserstein).",
            "Datasheet for Datasets + Model Card por modelo registrado.",
            "Testes E2E via docker compose + pytest-compose.",
            "Avaliacao de fairness (Disparate Impact, Equal Opportunity).",
            "",
            "# Medio prazo (cloud)",
            "Migrar Postgres para Cloud SQL (gerenciado).",
            "Migrar lake para GCS bucket (mlflow --default-artifact-root gs://...).",
            "Substituir FastAPI por Cloud Run (autoscale + HTTPS).",
            "Airflow gerenciado: Cloud Composer.",
            "",
            "# Governanca",
            "OAuth no Airflow e MLflow (em vez de admin/admin).",
            "Identity-Aware Proxy (IAP) no GCP para acesso autenticado.",
        ],
    )

    # Slide final
    slide_final(prs)

    prs.save(OUT_PPTX)
    print(f"\n{'='*60}\nSalvo: {OUT_PPTX}\nTotal: {len(prs.slides)} slides\n{'='*60}")


if __name__ == "__main__":
    main()
