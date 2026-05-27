"""Apresentacao PPTX da Pratica MLOps - padrao FGV (mesmo layout das lectures E1-E4).

Reusa ppt_helpers.py do outputs e gera diagramas em matplotlib.
Salva tambem os .mmd source no repositorio para documentacao.
"""
from __future__ import annotations

import sys
from pathlib import Path

# Importa ppt_helpers (que tem cores/fontes/header/footer FGV)
sys.path.insert(0, "/sessions/busy-hopeful-clarke/mnt/outputs")

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ppt_helpers import (
    AMBAR, AZUL, AZUL_CLARO, AZUL_MUITO_CLARO, BRANCO,
    CINZA, CINZA_CLARO, CINZA_ESCURO, SLIDE_H, SLIDE_W,
    VERDE, VERMELHO,
    add_rect, add_text, add_textbox, add_text_run,
    new_pres, titulo_slide,
)

# -----------------------------------------------------------------------------
PROJECT_DIR = Path("/sessions/busy-hopeful-clarke/mnt/FGV/mlops-pratica")
DOCS_DIR = PROJECT_DIR / "docs"
IMG_DIR = DOCS_DIR / "img"
MMD_DIR = DOCS_DIR / "mermaid"
IMG_DIR.mkdir(parents=True, exist_ok=True)
MMD_DIR.mkdir(parents=True, exist_ok=True)

# Cores hex (para matplotlib)
HX_AZUL = "#1F4E79"
HX_AZUL_CLARO = "#D5E8F0"
HX_AZUL_MC = "#EAF3F8"
HX_VERMELHO = "#C00000"
HX_CINZA = "#595959"
HX_CINZA_ESC = "#3F3F3F"
HX_VERDE = "#5A8F4B"
HX_AMBAR = "#C28A00"
HX_BRANCO = "#FFFFFF"
HX_LARANJA = "#ED7D31"


# -----------------------------------------------------------------------------
# Diagramas matplotlib (estilo limpo, padrao FGV)
# -----------------------------------------------------------------------------
def mbox(ax, x, y, w, h, text, color=HX_AZUL, fontcolor="white", fontsize=10):
    p = FancyBboxPatch((x, y), w, h,
                       boxstyle="round,pad=0.005,rounding_size=0.02",
                       linewidth=1.0, edgecolor=color, facecolor=color)
    ax.add_patch(p)
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
            color=fontcolor, fontsize=fontsize, weight="bold")


def marrow(ax, xy1, xy2, color=HX_CINZA, style="->", lw=1.8):
    a = FancyArrowPatch(xy1, xy2, arrowstyle=style, color=color, linewidth=lw,
                        mutation_scale=18, shrinkA=4, shrinkB=4)
    ax.add_patch(a)


def setup_fig(title, w=13.5, h=6.5):
    fig, ax = plt.subplots(figsize=(w, h), dpi=170)
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 7)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title(title, fontsize=14, weight="bold", color=HX_AZUL, pad=8)
    return fig, ax


def save(fig, name):
    path = IMG_DIR / f"{name}.png"
    fig.savefig(path, bbox_inches="tight", facecolor="white", dpi=180)
    plt.close(fig)
    return str(path)


def diag_arquitetura():
    fig, ax = setup_fig("Arquitetura - 7 containers Docker")
    mbox(ax, 0.3, 5.5, 2.2, 0.9, "HackerNews\nAPI", color=HX_VERDE)
    mbox(ax, 3.1, 5.5, 2.7, 0.9, "Airflow\nweb + scheduler", color=HX_AZUL)
    mbox(ax, 6.4, 5.5, 2.5, 0.9, "MinIO\ndata lake S3", color=HX_AZUL)
    mbox(ax, 9.6, 5.5, 3.0, 0.9, "MLflow\nTracking+Registry", color=HX_AZUL)
    mbox(ax, 5.0, 3.5, 4.0, 0.9, "Postgres + pgvector\nairflow / mlflow / app", color=HX_LARANJA)
    mbox(ax, 9.6, 3.5, 3.0, 0.9, "FastAPI\n/predict /search", color=HX_AZUL)
    mbox(ax, 3.1, 1.5, 2.7, 0.9, "ML Preditivo\nsklearn + XGBoost", color=HX_CINZA)
    mbox(ax, 6.4, 1.5, 2.5, 0.9, "Embeddings\nMiniLM-L6 (CPU)", color=HX_CINZA)

    marrow(ax, (2.5, 5.95), (3.1, 5.95))
    marrow(ax, (5.8, 5.95), (6.4, 5.95))
    marrow(ax, (8.9, 5.95), (9.6, 5.95))
    marrow(ax, (4.45, 5.5), (4.45, 2.4))
    marrow(ax, (7.65, 5.5), (7.65, 2.4))
    marrow(ax, (5.8, 2.0), (5.0, 3.5), color=HX_LARANJA)
    marrow(ax, (7.65, 2.4), (7.65, 3.5), color=HX_LARANJA)
    marrow(ax, (11.1, 5.5), (11.1, 4.4))
    return save(fig, "01_arquitetura")


def diag_eng_dados():
    fig, ax = setup_fig("Pipeline de Engenharia de Dados (Airflow @hourly)")
    y = 4.7
    mbox(ax, 0.3, y, 2.5, 1.0, "1. extract\nHN topstories", HX_VERDE)
    mbox(ax, 3.1, y, 2.5, 1.0, "2. raw\nMinIO parquet", HX_AZUL)
    mbox(ax, 5.9, y, 2.5, 1.0, "3. curate\ndedup por id", HX_AZUL)
    mbox(ax, 8.7, y, 2.5, 1.0, "4a. features\ntabular", HX_LARANJA)
    mbox(ax, 8.7, y - 2.0, 2.5, 1.0, "4b. features\ntext", HX_LARANJA)
    mbox(ax, 11.5, y, 2.2, 1.0, "treino\npreditivo", HX_CINZA)
    mbox(ax, 11.5, y - 2.0, 2.2, 1.0, "treino\nembeddings", HX_CINZA)
    marrow(ax, (2.8, 5.2), (3.1, 5.2))
    marrow(ax, (5.6, 5.2), (5.9, 5.2))
    marrow(ax, (8.4, 5.2), (8.7, 5.2))
    marrow(ax, (7.15, 4.7), (9.95, 3.7))
    marrow(ax, (11.2, 5.2), (11.5, 5.2))
    marrow(ax, (11.2, 3.2), (11.5, 3.2))
    ax.text(0.3, 6.5, "schedule: @hourly", fontsize=11, color=HX_AZUL, weight="bold",
            bbox=dict(facecolor=HX_AZUL_MC, edgecolor=HX_AZUL, pad=4))
    return save(fig, "02_eng_dados")


def diag_mlops():
    fig, ax = setup_fig("Fluxo MLOps - Tracking, Registry e Serving (MLflow)")
    mbox(ax, 0.3, 5.0, 2.5, 1.0, "load features\nMinIO", HX_AZUL)
    mbox(ax, 0.3, 3.0, 2.5, 1.0, "train\nRF + XGBoost", HX_AZUL)
    mbox(ax, 0.3, 1.0, 2.5, 1.0, "evaluate\nAUC/F1/PR-AUC", HX_CINZA)
    mbox(ax, 3.9, 3.8, 3.0, 2.2,
         "MLflow Tracking\nparams + metrics\nartifacts\nsignature", HX_LARANJA, fontsize=10)
    mbox(ax, 8.0, 4.6, 3.0, 1.0, "Registry\nNone -> Staging", HX_AZUL)
    mbox(ax, 8.0, 2.8, 3.0, 1.0, "Promote\n-> Production", HX_VERDE)
    mbox(ax, 11.5, 3.5, 2.2, 1.5, "FastAPI\n/predict\n/search", HX_AZUL)
    marrow(ax, (1.55, 5.0), (1.55, 4.0))
    marrow(ax, (1.55, 3.0), (1.55, 2.0))
    marrow(ax, (2.8, 5.5), (3.9, 5.0))
    marrow(ax, (2.8, 3.5), (3.9, 5.0))
    marrow(ax, (6.9, 5.0), (8.0, 5.0))
    marrow(ax, (9.5, 4.6), (9.5, 3.8))
    marrow(ax, (11.0, 3.3), (11.5, 4.25))
    return save(fig, "03_mlops")


def diag_devops():
    fig, ax = setup_fig("Fluxo DevOps - Gitflow + CI/CD")
    ax.plot([0.5, 13.5], [5.5, 5.5], color=HX_AZUL, lw=3)
    ax.plot([0.5, 13.5], [4.0, 4.0], color=HX_LARANJA, lw=3)
    ax.plot([2.0, 6.0], [2.5, 2.5], color=HX_VERDE, lw=3)
    ax.plot([7.0, 10.0], [2.5, 2.5], color=HX_CINZA, lw=3)
    ax.text(0.0, 5.55, "main", color=HX_AZUL, weight="bold", fontsize=11, ha="right")
    ax.text(0.0, 4.05, "develop", color=HX_LARANJA, weight="bold", fontsize=11, ha="right")
    ax.text(0.0, 2.55, "feature/* | release/*", color=HX_VERDE, weight="bold", fontsize=10, ha="right")
    for x in [2.0, 5.0, 9.5, 13.0]:
        ax.add_patch(plt.Circle((x, 5.5), 0.13, color=HX_AZUL))
    for x in [1.0, 3.5, 6.0, 8.5, 11.0, 12.5]:
        ax.add_patch(plt.Circle((x, 4.0), 0.13, color=HX_LARANJA))
    for x in [3.0, 4.5, 5.5]:
        ax.add_patch(plt.Circle((x, 2.5), 0.13, color=HX_VERDE))
    for x in [7.5, 8.5, 9.5]:
        ax.add_patch(plt.Circle((x, 2.5), 0.13, color=HX_CINZA))
    marrow(ax, (5.5, 2.5), (6.0, 4.0), color=HX_VERDE)
    marrow(ax, (9.5, 2.5), (9.5, 4.0), color=HX_CINZA)
    marrow(ax, (9.5, 4.0), (9.5, 5.5), color=HX_CINZA, lw=2)
    mbox(ax, 0.5, 0.3, 3.0, 0.9, "CI feature\nruff + pytest", HX_VERDE, fontsize=9)
    mbox(ax, 4.0, 0.3, 3.0, 0.9, "CI release\ndocker build", HX_CINZA, fontsize=9)
    mbox(ax, 7.5, 0.3, 3.0, 0.9, "CD staging\ndeploy auto", HX_AZUL, fontsize=9)
    mbox(ax, 11.0, 0.3, 2.7, 0.9, "CD prod\ntag v* manual", HX_LARANJA, fontsize=9)
    return save(fig, "04_devops")


def diag_dados():
    fig, ax = setup_fig("Base de Dados - HN API e camadas do lake")
    mbox(ax, 0.3, 5.5, 2.5, 1.0, "HN API\nfirebaseio", HX_VERDE)
    ax.text(0.3, 5.1, "/topstories.json\n/item/{id}.json", fontsize=8, color=HX_CINZA_ESC)
    mbox(ax, 3.5, 5.5, 2.5, 1.0, "raw\nparquet MinIO", HX_AZUL)
    ax.text(3.5, 5.1, "raw/hn/dt=YYYY-MM-DD/hr=HH/", fontsize=8, color=HX_CINZA_ESC)
    mbox(ax, 6.7, 5.5, 2.5, 1.0, "curated\nparquet MinIO", HX_AZUL)
    ax.text(6.7, 5.1, "dedup por id (score max)", fontsize=8, color=HX_CINZA_ESC)
    mbox(ax, 9.9, 5.5, 2.5, 1.0, "features\nparquet MinIO", HX_LARANJA)
    ax.text(9.9, 5.1, "tabular + text", fontsize=8, color=HX_CINZA_ESC)
    mbox(ax, 3.5, 2.2, 9.0, 1.5,
         "Postgres (com pgvector)\nDBs: airflow | mlflow | app (embeddings vector(384))",
         HX_LARANJA, fontsize=11)
    marrow(ax, (2.8, 6.0), (3.5, 6.0))
    marrow(ax, (6.0, 6.0), (6.7, 6.0))
    marrow(ax, (9.2, 6.0), (9.9, 6.0))
    marrow(ax, (11.15, 5.5), (11.15, 3.7), color=HX_LARANJA)
    return save(fig, "05_dados")


def diag_gcp():
    fig, ax = setup_fig("Deploy em GCP - VM Compute Engine + Persistent SSD")
    mbox(ax, 0.3, 5.5, 2.0, 1.0, "Internet", HX_CINZA, fontsize=11)
    mbox(ax, 3.0, 5.5, 3.0, 1.0, "Firewall (tag mlops)\n8080/5000/9001/8000", HX_LARANJA, fontsize=10)
    mbox(ax, 6.7, 4.0, 7.0, 2.5,
         "Compute Engine\ne2-standard-4, Ubuntu 22.04\nDocker + Compose: 7 containers",
         HX_AZUL, fontsize=11)
    mbox(ax, 6.7, 1.5, 4.5, 1.5,
         "Persistent SSD 200GB\n/opt/mlops/volumes/{postgres,minio}",
         HX_CINZA, fontsize=10)
    mbox(ax, 11.5, 1.7, 2.2, 0.6, "Cloud SQL (futuro)", HX_VERDE, fontsize=9)
    mbox(ax, 11.5, 0.9, 2.2, 0.6, "GCS bucket (futuro)", HX_VERDE, fontsize=9)
    marrow(ax, (2.3, 6.0), (3.0, 6.0))
    marrow(ax, (6.0, 6.0), (6.7, 6.0))
    marrow(ax, (10.2, 4.0), (8.95, 3.0))
    return save(fig, "06_gcp")


# -----------------------------------------------------------------------------
MERMAID = {
    "01_arquitetura.mmd": """flowchart LR
  API[HackerNews API] --> AF[Airflow]
  AF --> MIN[(MinIO data lake)]
  MIN --> MLF[MLflow Tracking+Registry]
  MIN --> M1[Preditivo sklearn+XGBoost]
  MIN --> M2[Embeddings MiniLM]
  M1 --> MLF
  M2 --> MLF
  M2 --> PG[(Postgres pgvector)]
  MLF --> SRV[FastAPI /predict /search]
  PG --> SRV
""",
    "02_eng_dados.mmd": """flowchart LR
  A[extract HN] --> B[raw/hn/dt=YYYY-MM-DD/hr=HH]
  B --> C[curate dedup]
  C --> D1[features tabular]
  C --> D2[features text]
  D1 --> E1[treino preditivo]
  D2 --> E2[treino embeddings]
""",
    "03_mlops.mmd": """flowchart LR
  L[Load features] --> T[Train RF + XGBoost]
  T --> E[Evaluate AUC/F1/PR-AUC]
  E --> M[MLflow Tracking]
  M --> R[Registry: Staging]
  R --> P[Production]
  P --> API[FastAPI /predict]
""",
    "04_devops.mmd": """gitGraph
  commit id: "init"
  branch develop
  checkout develop
  commit id: "scaffold"
  branch feature/ingestion
  checkout feature/ingestion
  commit id: "hn_client"
  commit id: "extractor"
  checkout develop
  merge feature/ingestion
  branch release/v0.1
  checkout release/v0.1
  commit id: "tag v0.1"
  checkout main
  merge release/v0.1
""",
    "05_dados.mmd": """flowchart LR
  API[HN API] --> RAW[raw parquet MinIO]
  RAW --> CUR[curated parquet]
  CUR --> FT[features tabular]
  CUR --> FX[features text]
  FT --> PG[(Postgres pgvector)]
  FX --> PG
""",
    "06_gcp.mmd": """flowchart LR
  INET[Internet] --> FW[Firewall 8080/5000/9001/8000]
  FW --> VM[Compute Engine e2-standard-4]
  VM --> DISK[(Persistent SSD 200GB)]
  VM -.-> CSQL[(Cloud SQL futuro)]
  VM -.-> GCS[(GCS bucket futuro)]
""",
}


def write_mermaid():
    for fname, src in MERMAID.items():
        p = MMD_DIR / fname
        p.write_text(src, encoding="utf-8")


# -----------------------------------------------------------------------------
# Slides (padrao FGV)
# -----------------------------------------------------------------------------
def header_pratica(slide, num, total, secao=""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.45), fill=AZUL)
    add_text(slide, Inches(0.4), Inches(0.05), Inches(10), Inches(0.35),
             "FGV — MBA IA & Analytics  |  MLOps  |  Pratica Integradora" +
             (f"  |  {secao}" if secao else ""),
             size=10, bold=True, color=BRANCO)
    add_text(slide, Inches(12.0), Inches(0.05), Inches(1.2), Inches(0.35),
             f"{num}/{total}", size=10, color=BRANCO, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.03), fill=AZUL)
    add_text(slide, Inches(0.4), Inches(7.30), Inches(10), Inches(0.18),
             "Prof. Andre Insardi  |  ext.andre.insardi@prof.fgv.edu.br",
             size=9, color=CINZA)
    add_text(slide, Inches(11.5), Inches(7.30), Inches(1.6), Inches(0.18),
             "MLOps - 2026", size=9, color=CINZA, align=PP_ALIGN.RIGHT)


def slide_capa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, Inches(4.7), SLIDE_H, fill=AZUL)
    add_text(s, Inches(0.5), Inches(0.5), Inches(3.7), Inches(0.4),
             "FGV - EDUCACAO EXECUTIVA", size=12, bold=True, color=BRANCO)
    add_text(s, Inches(0.5), Inches(1.2), Inches(3.7), Inches(0.4),
             "MBA em Inteligencia Artificial e", size=13, color=BRANCO)
    add_text(s, Inches(0.5), Inches(1.5), Inches(3.7), Inches(0.4),
             "Analytics Aplicadas a Negocios", size=13, color=BRANCO)
    add_text(s, Inches(0.5), Inches(6.4), Inches(3.7), Inches(0.4),
             "Prof. Andre Insardi", size=12, bold=True, color=BRANCO)
    add_text(s, Inches(0.5), Inches(6.7), Inches(3.7), Inches(0.4),
             "ext.andre.insardi@prof.fgv.edu.br", size=10, color=BRANCO)
    add_text(s, Inches(0.5), Inches(7.0), Inches(3.7), Inches(0.4),
             "Maio/2026", size=10, color=BRANCO)
    add_text(s, Inches(5.2), Inches(1.6), Inches(8.0), Inches(0.5),
             "PRATICA INTEGRADORA", size=18, bold=True, color=VERMELHO)
    add_text(s, Inches(5.2), Inches(2.1), Inches(8.0), Inches(0.5),
             "MLOps end-to-end | 100% Docker", size=14, color=CINZA)
    add_text(s, Inches(5.2), Inches(2.9), Inches(8.0), Inches(1.2),
             "HackerNews MLOps", size=44, bold=True, color=AZUL)
    add_text(s, Inches(5.2), Inches(4.4), Inches(8.0), Inches(0.5),
             "Airflow + MLflow + pgvector + FastAPI", size=18, color=CINZA)
    add_text(s, Inches(5.2), Inches(4.8), Inches(8.0), Inches(0.5),
             "Deploy local e em VM GCP", size=18, color=CINZA)
    add_rect(s, Inches(5.2), Inches(5.8), Inches(7.5), Inches(0.04), fill=AZUL)
    add_text(s, Inches(5.2), Inches(6.0), Inches(8.0), Inches(0.4),
             "Pratica integradora dos encontros E1-E4: codigo, dados e modelos versionados.",
             size=11, italic=True, color=CINZA)
    return s


def slide_secao_divider(prs, num, total, num_secao, titulo, descricao):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_pratica(s, num, total, f"Bloco {num_secao}")
    add_rect(s, 0, Inches(0.45), SLIDE_W, Inches(6.8), fill=AZUL_MUITO_CLARO)
    add_text(s, Inches(0.6), Inches(1.8), Inches(2.5), Inches(2.0),
             num_secao, size=120, bold=True, color=AZUL, align=PP_ALIGN.LEFT)
    add_text(s, Inches(3.2), Inches(2.5), Inches(10.0), Inches(0.8),
             titulo, size=36, bold=True, color=AZUL)
    add_rect(s, Inches(3.2), Inches(3.55), Inches(2.0), Inches(0.05), fill=VERMELHO)
    add_text(s, Inches(3.2), Inches(3.8), Inches(10.0), Inches(2.0),
             descricao, size=14, color=CINZA_ESCURO)
    return s


def slide_bullets(prs, num, total, secao, titulo, subtitulo, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_pratica(s, num, total, secao)
    titulo_slide(s, titulo, subtitulo)
    tb, tf = add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0))
    for i, b in enumerate(bullets):
        if b.startswith("# "):
            add_text_run(tf, b[2:], size=15, bold=True, color=AZUL, first=(i == 0))
        elif b == "":
            add_text_run(tf, "", size=8, first=(i == 0))
        else:
            add_text_run(tf, "- " + b, size=13, color=CINZA_ESCURO, first=(i == 0))
    return s


def slide_imagem(prs, num, total, secao, titulo, subtitulo, img_path, caption=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_pratica(s, num, total, secao)
    titulo_slide(s, titulo, subtitulo)
    s.shapes.add_picture(img_path, Inches(0.6), Inches(1.9), width=Inches(12.1))
    if caption:
        add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
                 caption, size=11, italic=True, color=CINZA, align=PP_ALIGN.CENTER)
    return s


def slide_stack_tabela(prs, num, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_pratica(s, num, total, "Stack")
    titulo_slide(s, "Stack confirmada", "7 containers Docker, 100% offline-capable")
    linhas = [
        ("Orquestracao", "Apache Airflow 2.9 (LocalExecutor)"),
        ("Tracking + Registry", "MLflow 2.16 (backend Postgres, artifact MinIO)"),
        ("Data Lake", "MinIO (S3-compatible)"),
        ("Metadata + Vector", "Postgres 16 + pgvector 0.7 (HNSW cosine)"),
        ("Embeddings", "sentence-transformers/all-MiniLM-L6-v2 (offline, dim=384)"),
        ("ML Preditivo", "scikit-learn 1.5 + XGBoost 2.1"),
        ("Serving", "FastAPI 0.115 + Uvicorn"),
        ("CI/CD", "GitHub Actions (ruff + pytest + docker build)"),
    ]
    shape = s.shapes.add_table(len(linhas) + 1, 2, Inches(0.5), Inches(2.0),
                               Inches(12.3), Inches(4.8))
    table = shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.8)
    table.cell(0, 0).text = "Camada"
    table.cell(0, 1).text = "Tecnologia"
    for c in [0, 1]:
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = AZUL
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12); r.font.bold = True
                r.font.color.rgb = BRANCO; r.font.name = "Calibri"
    for i, (camada, tec) in enumerate(linhas, 1):
        table.cell(i, 0).text = camada
        table.cell(i, 1).text = tec
        cor = CINZA_CLARO if i % 2 == 0 else BRANCO
        for c in [0, 1]:
            cell = table.cell(i, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = cor
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11); r.font.color.rgb = CINZA_ESCURO
                    r.font.name = "Calibri"
                    if c == 0:
                        r.font.bold = True
    return s


def slide_obrigado(prs, num, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=AZUL)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12.0), Inches(1.5),
             "Obrigado!", size=80, bold=True, color=BRANCO)
    add_rect(s, Inches(0.6), Inches(4.2), Inches(2.0), Inches(0.05), fill=BRANCO)
    add_text(s, Inches(0.6), Inches(4.4), Inches(12.0), Inches(0.6),
             "Perguntas, criticas e sugestoes sao bem-vindas.",
             size=20, color=BRANCO)
    add_text(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(0.5),
             "Prof. Andre Insardi  •  ext.andre.insardi@prof.fgv.edu.br",
             size=14, italic=True, color=BRANCO)
    add_text(s, Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.4),
             "FGV - MBA IA & Analytics  |  MLOps  |  Maio/2026",
             size=11, color=BRANCO)
    return s


# -----------------------------------------------------------------------------
def main():
    print("==> Gerando .mmd files")
    write_mermaid()

    print("==> Gerando PNGs dos diagramas")
    img_arq = diag_arquitetura()
    img_eng = diag_eng_dados()
    img_mlo = diag_mlops()
    img_dev = diag_devops()
    img_dad = diag_dados()
    img_gcp = diag_gcp()

    print("==> Montando PPTX (padrao FGV)")
    prs = new_pres()

    slides_plan = [
        ("capa", None),
        ("agenda", None),

        ("secao", ("1", "Problema",
                   "Por que precisamos de uma pratica integradora end-to-end e qual e o caso de uso escolhido.")),
        ("problema1", None),
        ("problema2", None),

        ("secao", ("2", "Base de dados",
                   "Fonte, camadas e modelagem - HackerNews API + MinIO + Postgres/pgvector.")),
        ("dados_text", None),
        ("dados_img", img_dad),

        ("secao", ("3", "Arquitetura geral",
                   "Sete containers Docker em uma rede unica, sem dependencias de cloud.")),
        ("stack", None),
        ("arq_img", img_arq),

        ("secao", ("4", "Engenharia de Dados (Airflow)",
                   "DAGs, particionamento, idempotencia e features.")),
        ("eng_img", img_eng),
        ("eng_design", None),

        ("secao", ("5", "Fluxo MLOps (MLflow)",
                   "Tracking, Registry, promocao para Production e serving via FastAPI.")),
        ("mlops_img", img_mlo),
        ("mlops_3v", None),
        ("mlops_pred", None),
        ("mlops_emb", None),

        ("secao", ("6", "Fluxo DevOps (Gitflow)",
                   "Branches, CI no GitHub Actions e CD por estagio.")),
        ("devops_img", img_dev),
        ("devops_ci", None),

        ("secao", ("7", "Deploy em VM GCP",
                   "Compute Engine + Persistent Disk + Firewall, via gcloud script.")),
        ("gcp_img", img_gcp),
        ("gcp_how", None),

        ("secao", ("8", "Conclusoes",
                   "O que entregamos e proximos passos para evolucao.")),
        ("concl1", None),
        ("concl2", None),
        ("obrigado", None),
    ]
    total = len(slides_plan)

    n = 0
    for tipo, arg in slides_plan:
        n += 1
        if tipo == "capa":
            slide_capa(prs)
        elif tipo == "agenda":
            slide_bullets(prs, n, total, "Agenda",
                "Agenda da Pratica Integradora",
                "Roteiro de 8 blocos do problema ao deploy em GCP",
                [
                    "# Estrutura (8 blocos)",
                    "1. Problema e caso de uso (HackerNews 'vai bombar')",
                    "2. Base de dados (API + camadas + Postgres/pgvector)",
                    "3. Arquitetura geral (7 containers Docker)",
                    "4. Pipeline de Engenharia de Dados (Airflow)",
                    "5. Fluxo MLOps (MLflow Tracking + Registry + Serving)",
                    "6. Fluxo DevOps (Gitflow + CI/CD GitHub Actions)",
                    "7. Deploy em VM GCP (Compute Engine + SSD)",
                    "8. Conclusoes e proximos passos",
                    "",
                    "# Materiais",
                    "Repositorio: mlops-pratica/ (codigo + DAGs + docs)",
                    "Diagramas Mermaid: docs/mermaid/*.mmd",
                    "Guia GCP: docs/deploy_gcp_vm.md",
                ])
        elif tipo == "secao":
            num_secao, titulo, descricao = arg
            slide_secao_divider(prs, n, total, num_secao, titulo, descricao)
        elif tipo == "problema1":
            slide_bullets(prs, n, total, "Problema",
                "O que esta pratica resolve",
                "Materializar toda a cadeia MLOps em um exemplo rodavel em laptop",
                [
                    "# Motivacao",
                    "Cada encontro (E1-E4) discutiu conceitos isolados. Falta integrar.",
                    "Os alunos precisam VER o ciclo completo funcionando end-to-end.",
                    "",
                    "# Escopo",
                    "Ingestao real de API publica (HackerNews) com Airflow @hourly.",
                    "Dois pipes de ML lado a lado: tabular (classificador) e vetorial (embeddings).",
                    "Versionamento triplo: codigo (Git), dados (lake), modelos (MLflow).",
                    "Serving via FastAPI carregando modelo do Registry.",
                    "",
                    "# Restricoes",
                    "100% Docker - sem dependencia de cloud durante o desenvolvimento.",
                    "Modelos offline-capable (CPU): MiniLM + sklearn + XGBoost.",
                    "Reprodutivel via 'make up' (laptop) e em VM GCP (mesmo compose).",
                ])
        elif tipo == "problema2":
            slide_bullets(prs, n, total, "Caso de uso",
                "Caso de uso: prever se um post HN vai bombar",
                "Target binario derivado do score real do HackerNews",
                [
                    "# Definicao do target",
                    "vai_bombar = (score >= 100), threshold didatico (gera ~10-20% positivos).",
                    "",
                    "# Features tabulares",
                    "Comprimento e numero de palavras do titulo.",
                    "Hora do post e dia da semana (efeito de horario nobre).",
                    "Presenca de URL externa, dominio (one-hot com min_frequency=5).",
                    "Autor (encoded), termina com '?' (Ask HN).",
                    "",
                    "# Pipe vetorial complementar",
                    "Embeddings de titulos para busca semantica (RAG simples).",
                    "Aplicacao: deduplicar topicos, achar posts similares.",
                    "",
                    "# Aplicacoes empresariais analogas",
                    "Priorizar feed editorial; detectar conteudo viral; classificar tickets.",
                ])
        elif tipo == "dados_text":
            slide_bullets(prs, n, total, "Base de dados",
                "Fonte e camadas de dados",
                "API publica + lake parquet + Postgres com pgvector",
                [
                    "# Fonte",
                    "HackerNews Firebase API (sem auth, sem rate limit pratico).",
                    "Endpoints: /v0/topstories.json (IDs) e /v0/item/{id}.json (detalhe).",
                    "",
                    "# Lake (MinIO, formato parquet)",
                    "raw/hn/dt=YYYY-MM-DD/hr=HH/items.parquet",
                    "curated/hn/stories.parquet (dedup por id, score maximo)",
                    "features/tabular/train.parquet (input do pipe preditivo)",
                    "features/text/titles.parquet (input do pipe de embeddings)",
                    "",
                    "# Postgres (3 databases na mesma instancia)",
                    "airflow - metadata do scheduler",
                    "mlflow - backend store do Tracking Server (runs + Registry)",
                    "app - tabela embeddings(id, title, vector(384), score, ts) + indice HNSW",
                ])
        elif tipo == "dados_img":
            slide_imagem(prs, n, total, "Base de dados",
                "Visao das camadas", "Da API ao Postgres via MinIO parquet",
                arg, caption="Camadas raw, curated e features no MinIO + Postgres com pgvector.")
        elif tipo == "stack":
            slide_stack_tabela(prs, n, total)
        elif tipo == "arq_img":
            slide_imagem(prs, n, total, "Arquitetura",
                "Arquitetura geral", "Sete containers em uma rede Docker unica (mlops-net)",
                arg, caption="Postgres + MinIO + MLflow + Airflow (web+sched) + FastAPI.")
        elif tipo == "eng_img":
            slide_imagem(prs, n, total, "Eng. Dados",
                "DAG pipeline_ingestao", "Schedule @hourly - 4 etapas + 2 outputs paralelos",
                arg, caption="extract -> raw -> curated -> features (tabular + text).")
        elif tipo == "eng_design":
            slide_bullets(prs, n, total, "Eng. Dados",
                "Padroes de design das DAGs", "Decisoes que mantem o pipeline robusto",
                [
                    "# TaskFlow API",
                    "Decorators @dag e @task (Python puro, sem operadores legados).",
                    "XCom automatico via return + assinatura entre tasks.",
                    "",
                    "# Idempotencia e particionamento",
                    "Cada execucao sobrescreve a particao da hora corrente.",
                    "Particionamento Hive-style (dt=, hr=) habilita query eficiente futura.",
                    "",
                    "# Dedup na curated",
                    "Concat (curated + raw novo) -> sort score desc -> drop_duplicates(id).",
                    "Snapshot estavel sem perder o melhor score ja visto.",
                    "",
                    "# Resilience",
                    "Cliente HTTP com tenacity (3 retries, backoff exponencial).",
                    "ThreadPoolExecutor 16 workers para hidratar 100 items em ~5s.",
                    "Retry da task (2x, 2min) cobre falhas transientes da API.",
                ])
        elif tipo == "mlops_img":
            slide_imagem(prs, n, total, "MLOps",
                "Fluxo MLOps completo", "Tracking, Registry e Serving",
                arg, caption="Dois experimentos -> Registry -> FastAPI /predict.")
        elif tipo == "mlops_3v":
            slide_bullets(prs, n, total, "MLOps",
                "Triplo versionamento", "Codigo, dados e modelos versionados em ferramentas dedicadas",
                [
                    "# 1. Codigo - Git",
                    "Branches: main (estavel), develop (integracao), feature/* e release/*.",
                    "Tags v* para releases (v0.1, v0.2, ...).",
                    "",
                    "# 2. Dados - Parquet particionado em MinIO",
                    "raw/ guarda o estado bruto da hora (auditavel).",
                    "curated/ e a versao estavel (snapshot).",
                    "features/ versiona o dataset de treino (pode evoluir para DVC).",
                    "",
                    "# 3. Modelos - MLflow Tracking + Registry",
                    "Cada run captura: params, metrics, artifacts (pipeline sklearn completo).",
                    "Signature inferida automaticamente.",
                    "Registry com versoes numericas e stages (None|Staging|Production|Archived).",
                ])
        elif tipo == "mlops_pred":
            slide_bullets(prs, n, total, "MLOps",
                "DAG treino preditivo (02:00 daily)",
                "Compara RF e XGBoost, promove vencedor por ROC-AUC",
                [
                    "# Etapas",
                    "load_features -> train_and_register_models -> promote.",
                    "",
                    "# Modelos comparados",
                    "RandomForest (n_estimators=200, max_depth=12).",
                    "XGBoost (n_estimators=300, max_depth=6, learning_rate=0.1).",
                    "",
                    "# Metricas logadas",
                    "accuracy, f1, roc_auc, pr_auc em test split estratificado 20%.",
                    "Hiperparametros via mlflow.log_params.",
                    "",
                    "# Politica de promocao",
                    "DAG promove ultima versao para Staging automaticamente.",
                    "Producao real: gate humano + comparativo de metricas no Registry.",
                ])
        elif tipo == "mlops_emb":
            slide_bullets(prs, n, total, "MLOps",
                "DAG treino embeddings (03:00 daily)", "Pipe vetorial complementar com pgvector",
                [
                    "# Etapas",
                    "load_text -> encode -> index.",
                    "",
                    "# Encoder",
                    "sentence-transformers/all-MiniLM-L6-v2 (offline, CPU, dim=384).",
                    "Pre-baixado no Dockerfile (evita download em runtime).",
                    "Normalizacao L2 dos vetores (cosine distance correto).",
                    "",
                    "# Indexacao",
                    "Upsert em embeddings(id, title, vector, score, by_author, ts).",
                    "Indice HNSW vector_cosine_ops (rapido para milhoes de vetores).",
                    "",
                    "# Tracking no MLflow",
                    "Loga: modelo, dim, n_inputs, amostras (samples.json como artifact).",
                ])
        elif tipo == "devops_img":
            slide_imagem(prs, n, total, "DevOps",
                "Gitflow + CI/CD", "Branches main/develop/feature/release com gates por estagio",
                arg, caption="GitHub Actions roda ruff + pytest em PRs e docker build em release.")
        elif tipo == "devops_ci":
            slide_bullets(prs, n, total, "DevOps",
                "Workflow ci.yml (GitHub Actions)", "Dois jobs em sequencia: lint+test e docker build",
                [
                    "# Triggers",
                    "push em main/develop e pull_request em main/develop.",
                    "",
                    "# Job 1 - lint-test",
                    "Setup Python 3.11 + pip cache.",
                    "Install requirements-dev.txt.",
                    "ruff check src tests dags.",
                    "pytest --cov=src/mlops_pratica --cov-report=term-missing.",
                    "",
                    "# Job 2 - docker-build (needs lint-test)",
                    "docker buildx setup.",
                    "Build sanity-check das 3 imagens: airflow, mlflow, fastapi.",
                    "",
                    "# Politica de release",
                    "Tag vX.Y na main dispara workflow de release (pode habilitar push para GHCR).",
                    "Gate humano antes de qualquer deploy compartilhado.",
                ])
        elif tipo == "gcp_img":
            slide_imagem(prs, n, total, "Deploy GCP",
                "Arquitetura GCP (VM unica)",
                "Compute Engine e2-standard-4 + Persistent SSD para volumes",
                arg, caption="Firewall por tag mlops libera as 4 portas das UIs.")
        elif tipo == "gcp_how":
            slide_bullets(prs, n, total, "Deploy GCP",
                "Como provisionar", "Script unico (gcloud) + startup-script automatizado",
                [
                    "# Pre-requisitos",
                    "Projeto GCP com billing.",
                    "gcloud autenticado e cota para 4 vCPU.",
                    "",
                    "# Comando unico",
                    "export GCP_PROJECT=meu-projeto",
                    "./infra/gcp/provision.sh",
                    "",
                    "# O script faz",
                    "Cria disco SSD 200GB (pd-ssd).",
                    "Cria VM ubuntu-22 com startup-script.",
                    "Abre firewall por tag (ajuste SOURCE_IP_RANGE em prod!).",
                    "",
                    "# Startup-script automatiza",
                    "Instala Docker Engine + Compose plugin.",
                    "Formata e monta disco em /opt/mlops.",
                    "Clona o repositorio e roda docker compose up -d.",
                    "",
                    "# Custo estimado",
                    "~US$ 150/mes (VM 24/7 + 300GB SSD).",
                    "Preemptible reduz ~70%.",
                ])
        elif tipo == "concl1":
            slide_bullets(prs, n, total, "Conclusoes",
                "O que entregamos", "Stack MLOps end-to-end, reprodutivel e portavel",
                [
                    "# Componentes implementados",
                    "7 containers Docker orquestrados via docker-compose.",
                    "3 DAGs Airflow cobrindo ingestao, treino preditivo e embeddings.",
                    "MLflow Tracking + Registry com promocao automatizada.",
                    "Vector search em pgvector com indice HNSW cosine.",
                    "FastAPI servindo /predict (Registry) e /search (vetorial).",
                    "Triplo versionamento: Git + lake parquet + MLflow Registry.",
                    "CI completo (ruff + pytest + docker build) via GitHub Actions.",
                    "Manual de deploy em VM GCP com script de provisionamento.",
                    "",
                    "# Tudo reproduzivel",
                    "make up sobe a stack inteira em <10 minutos no laptop.",
                    "Mesmo compose roda na VM GCP via override de volumes.",
                ])
        elif tipo == "concl2":
            slide_bullets(prs, n, total, "Conclusoes",
                "Proximos passos (evolucao)", "Caminhos para amadurecer a stack rumo a producao",
                [
                    "# Curto prazo (laboratorio)",
                    "Monitoracao de drift com Evidently AI (KS, PSI, Wasserstein).",
                    "Datasheet for Datasets + Model Card por modelo registrado.",
                    "Testes E2E via docker compose + pytest-compose.",
                    "Avaliacao de fairness (Disparate Impact, Equal Opportunity).",
                    "",
                    "# Medio prazo (cloud)",
                    "Migrar Postgres para Cloud SQL (gerenciado).",
                    "Migrar lake para GCS bucket (mlflow --default-artifact-root gs://...).",
                    "Substituir FastAPI por Cloud Run (autoscale + HTTPS).",
                    "Airflow gerenciado: Cloud Composer.",
                    "",
                    "# Governanca",
                    "OAuth no Airflow e MLflow (em vez de admin/admin).",
                    "Identity-Aware Proxy (IAP) no GCP para acesso autenticado.",
                    "Politicas de retention de runs e snapshots de disk.",
                ])
        elif tipo == "obrigado":
            slide_obrigado(prs, n, total)

    out = PROJECT_DIR / "Apresentacao_MLOps_Pratica.pptx"
    prs.save(out)
    print(f"\n{'='*60}")
    print(f"PPTX salvo: {out}")
    print(f"{len(prs.slides)} slides no padrao FGV.")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
